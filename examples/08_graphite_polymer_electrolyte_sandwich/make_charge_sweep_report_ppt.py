"""Build a strict Eg08.07 charge-sweep PowerPoint report.

The report intentionally does more than paste figures into slides.  It checks
whether all charge states really start from the same t=0 coordinates, rewrites
z-axis data onto a single CMC-facing graphite reference, redraws comparison
figures with consistent axes and colors, and writes a visual audit so blank
data slides are caught automatically.
"""

from __future__ import annotations

import csv
import hashlib
import json
import math
import os
import shutil
import subprocess
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Iterable

import numpy as np


DEFAULT_CASES = (
    ("-18 uC/cm2", "cmcface_m18p0_uC_cm2", -18.0),
    ("-9 uC/cm2", "cmcface_m9p0_uC_cm2", -9.0),
    ("-3 uC/cm2", "cmcface_m3p0_uC_cm2", -3.0),
    ("0 uC/cm2", "cmcface_00_uC_cm2", 0.0),
)

SPECIES_ORDER = ("EC", "EMC", "DEC", "LI", "PF6", "NA", "CMCNA")
PERMEANT_SPECIES = ("EC", "EMC", "DEC", "LI", "PF6")
CHARGE_PALETTE = {
    -18.0: "#2B5C8A",
    -9.0: "#2A9D8F",
    -3.0: "#E76F51",
    0.0: "#6D597A",
}
SPECIES_PALETTE = {
    "EC": "#4C78A8",
    "EMC": "#72B7B2",
    "DEC": "#F58518",
    "LI": "#D62728",
    "PF6": "#7F3C8D",
    "NA": "#54A24B",
    "CMCNA": "#8C6D31",
}
LI_SOLVATION_TARGETS = ("solvent_o", "cmc_o", "pf6_f")
LI_SOLVATION_PALETTE = {
    "solvent_o": "#3B82B8",
    "cmc_o": "#9A7B39",
    "pf6_f": "#7F3C8D",
}


@dataclass(frozen=True)
class CasePayload:
    label: str
    dirname: str
    charge: float
    analysis_dir: Path
    relax_dir: Path
    summary: dict[str, Any]
    z0_nm: float
    box_z_nm: float
    z_axis_direction: str
    z_axis_source: str
    z_axis_surface_charge_uC_cm2: float | None


@dataclass(frozen=True)
class ZAxisReference:
    z0_nm: float
    box_z_nm: float
    direction: str
    source: str
    surface_charge_uC_cm2: float | None
    layer_name: str | None = None
    region: str | None = None
    z_window_nm: tuple[float, float] | None = None


def _case_specs() -> tuple[tuple[str, str, float], ...]:
    """Return report cases, optionally overridden by JSON in the environment."""

    raw = os.environ.get("EG08_SWEEP_CASES_JSON")
    if not raw:
        return DEFAULT_CASES
    payload = json.loads(raw)
    cases: list[tuple[str, str, float]] = []
    for item in payload:
        cases.append((str(item["label"]), str(item["dir"]), float(item["charge_uC_cm2"])))
    return tuple(cases)


def _final_nvt_dir(item: CasePayload) -> Path:
    workflow = item.relax_dir / "05_relaxation_workflow"
    candidates = [
        workflow / "03_final_nvt_release",
        workflow / "04_final_nvt_release",
        workflow / "01_final_nvt",
        workflow / "final_nvt",
    ]
    for path in candidates:
        if (path / "md.xtc").is_file() and (path / "md.gro").is_file():
            return path
    if workflow.is_dir():
        for path in sorted(workflow.iterdir()):
            lname = path.name.lower()
            if path.is_dir() and "final" in lname and "nvt" in lname and (path / "md.xtc").is_file() and (path / "md.gro").is_file():
                return path
    return candidates[0]


def _read_json(path: Path) -> dict[str, Any]:
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
        return payload if isinstance(payload, dict) else {}
    except Exception:
        return {}


def _read_rows(path: Path) -> list[dict[str, str]]:
    if not path.is_file():
        return []
    with path.open("r", encoding="utf-8", errors="replace", newline="") as fh:
        return [dict(row) for row in csv.DictReader(fh)]


def _write_rows(path: Path, rows: Iterable[dict[str, Any]]) -> None:
    rows = [dict(row) for row in rows]
    path.parent.mkdir(parents=True, exist_ok=True)
    if not rows:
        path.write_text("", encoding="utf-8")
        return
    fields: list[str] = []
    for row in rows:
        for key in row:
            if key not in fields:
                fields.append(key)
    with path.open("w", encoding="utf-8", newline="") as fh:
        writer = csv.DictWriter(fh, fieldnames=fields)
        writer.writeheader()
        writer.writerows(rows)


def _safe_float(value: Any, default: float = np.nan) -> float:
    try:
        if value is None or value == "":
            return float(default)
        return float(value)
    except Exception:
        return float(default)


def _optional_float(value: Any) -> float | None:
    out = _safe_float(value, np.nan)
    return float(out) if np.isfinite(out) else None


def _box_from_gro(path: Path) -> tuple[float, float, float]:
    try:
        parts = path.read_text(encoding="utf-8", errors="replace").splitlines()[-1].split()
        if len(parts) >= 3:
            return float(parts[0]), float(parts[1]), float(parts[2])
    except Exception:
        pass
    return (1.0, 1.0, 1.0)


def _trapz(y: np.ndarray, x: np.ndarray) -> float:
    if hasattr(np, "trapezoid"):
        return float(np.trapezoid(y, x))
    if hasattr(np, "trapz"):
        return float(np.trapz(y, x))
    if y.size < 2 or x.size < 2:
        return 0.0
    return float(np.sum(0.5 * (y[1:] + y[:-1]) * (x[1:] - x[:-1])))


def _canonical_species(value: Any) -> str:
    text = str(value or "").strip().upper()
    if text.startswith("LI"):
        return "LI"
    if text.startswith("NA"):
        return "NA"
    if text.startswith("PF6"):
        return "PF6"
    if "CMC" in text or "POLY" in text:
        return "CMCNA"
    for key in ("EMC", "DEC", "EC"):
        if text == key or text.startswith(f"{key}_"):
            return key
    return text


def _species_matches(value: Any, target: str) -> bool:
    return _canonical_species(value) == _canonical_species(target)


def _configure_matplotlib() -> None:
    import matplotlib

    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    plt.rcParams.update(
        {
            "font.family": "Arial",
            "font.sans-serif": ["Arial", "Liberation Sans", "DejaVu Sans"],
            "axes.titlesize": 12,
            "axes.labelsize": 10,
            "xtick.labelsize": 8,
            "ytick.labelsize": 8,
            "legend.fontsize": 8,
            "svg.fonttype": "none",
            "pdf.fonttype": 42,
            "axes.spines.top": False,
            "axes.spines.right": False,
        }
    )


def _save_fig(fig, base_path: Path, *, dpi: int = 220) -> dict[str, Path]:
    png = base_path.with_suffix(".png")
    svg = base_path.with_suffix(".svg")
    png.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(svg, bbox_inches="tight")
    fig.savefig(png, dpi=dpi, bbox_inches="tight")
    return {"png": png, "svg": svg}


def _charge_color(charge: float) -> str:
    return CHARGE_PALETTE.get(float(charge), "#4D4D4D")


def _species_color(species: str) -> str:
    return SPECIES_PALETTE.get(_canonical_species(species), "#4D4D4D")


def _infer_z0_nm(summary: dict[str, Any]) -> float:
    phase_stats = dict(summary.get("phase_stats") or {})
    candidates = []
    for name, stats in phase_stats.items():
        lname = str(name).lower()
        if "graphite_bottom" in lname:
            z = stats.get("p95_z_nm") or stats.get("p100_z_nm") or stats.get("p99_z_nm")
            if z is not None:
                return float(z)
        if "graphite" in lname:
            z = stats.get("p95_z_nm") or stats.get("p100_z_nm") or stats.get("p99_z_nm")
            if z is not None:
                candidates.append(float(z))
    if candidates:
        return float(min(candidates))
    membrane = dict(summary.get("membrane_permeation") or {})
    if membrane.get("membrane_z_lo_nm") is not None:
        return float(membrane["membrane_z_lo_nm"])
    return 0.0


def _z_axis_from_patch(case_dir: Path, summary: dict[str, Any]) -> ZAxisReference:
    """Return the report z-axis reference.

    The charge-sweep convention is intentionally orientation-aware: x=0 is the
    CMC-facing graphite inner surface and +x points from that surface into the
    CMC/electrolyte stack.  For the current Eg08.07 geometry that surface is
    GRAPHITE_TOP/bottom, so the plotting direction is decreasing absolute z.
    """

    system_gro = case_dir / "02_system" / "system.gro"
    box_z = _box_from_gro(system_gro)[2] if system_gro.is_file() else np.nan
    patch = _read_json(case_dir / "02_system" / "charge_patch_report.json")
    regions = patch.get("regions") or patch.get("fixed_charge_regions") or []
    chosen: dict[str, Any] | None = None
    for region in regions:
        if str(region.get("label")) == "cmc_facing_graphite_inner_face":
            chosen = dict(region)
            break
    if chosen is None:
        for region in regions:
            layer = str(region.get("layer_name") or region.get("layer") or "").upper()
            face = str(region.get("region") or region.get("face") or "").lower()
            if layer == "GRAPHITE_TOP" and face == "bottom":
                chosen = dict(region)
                break
    if chosen is not None:
        window_raw = chosen.get("z_window_nm") or chosen.get("z_window") or []
        window = [_safe_float(v) for v in window_raw[:2]]
        window = [v for v in window if np.isfinite(v)]
        layer_name = str(chosen.get("layer_name") or chosen.get("layer") or "")
        region_name = str(chosen.get("region") or chosen.get("face") or "").lower()
        direction = "decreasing" if region_name == "bottom" else "increasing"
        if len(window) >= 2:
            zlo, zhi = min(window), max(window)
            z0 = zlo if direction == "decreasing" else zhi
            z_window = (zlo, zhi)
        elif len(window) == 1:
            z0 = window[0]
            z_window = (window[0], window[0])
        else:
            z0 = _infer_z0_nm(summary)
            z_window = None
        return ZAxisReference(
            z0_nm=float(z0),
            box_z_nm=float(box_z if np.isfinite(box_z) and box_z > 0 else 1.0),
            direction=direction,
            source="charge_patch_report:cmc_facing_graphite_inner_face",
            surface_charge_uC_cm2=_optional_float(chosen.get("surface_charge_uC_cm2")),
            layer_name=layer_name or None,
            region=region_name or None,
            z_window_nm=z_window,
        )
    return ZAxisReference(
        z0_nm=float(_infer_z0_nm(summary)),
        box_z_nm=float(box_z if np.isfinite(box_z) and box_z > 0 else 1.0),
        direction="decreasing",
        source="fallback:summary_phase_stats",
        surface_charge_uC_cm2=None,
    )


def _z_plot_nm(item: CasePayload, z_abs_nm: Any) -> float:
    z_abs = _safe_float(z_abs_nm)
    if not np.isfinite(z_abs):
        return float("nan")
    lz = max(float(item.box_z_nm), 1.0e-12)
    if item.z_axis_direction == "decreasing":
        return float((float(item.z0_nm) - z_abs) % lz)
    return float((z_abs - float(item.z0_nm)) % lz)


def _write_z_axis_reference(payloads: list[CasePayload], out_dir: Path) -> None:
    rows = []
    for item in payloads:
        rows.append(
            {
                "case": item.label,
                "charge_uC_cm2": item.charge,
                "z0_nm": item.z0_nm,
                "box_z_nm": item.box_z_nm,
                "direction": item.z_axis_direction,
                "source": item.z_axis_source,
                "surface_charge_uC_cm2_at_z0": item.z_axis_surface_charge_uC_cm2,
                "definition": "z_plot_nm=0 at CMC-facing graphite inner surface; +z_plot points into CMC, then electrolyte, then opposite graphite.",
            }
        )
    _write_rows(out_dir / "z_axis_reference.csv", rows)
    _write_json(
        out_dir / "z_axis_reference.json",
        {
            "definition": "z_plot_nm=0 at the negatively charged CMC-facing graphite inner surface; +z_plot points into CMC, then electrolyte, then opposite graphite, wrapped into one full periodic box.",
            "cases": rows,
        },
    )


def _load_payloads(root: Path, cases: tuple[tuple[str, str, float], ...]) -> list[CasePayload]:
    payloads: list[CasePayload] = []
    for label, dirname, charge in cases:
        relax_dir = root / dirname / "03_relaxation_sampling"
        analysis_dir = relax_dir / "06_analysis" / "layer_stack_interface"
        summary = _read_json(analysis_dir / "interface_profile_summary.json")
        zref = _z_axis_from_patch(root / dirname, summary)
        payloads.append(
            CasePayload(
                label=label,
                dirname=dirname,
                charge=float(charge),
                analysis_dir=analysis_dir,
                relax_dir=relax_dir,
                summary=summary,
                z0_nm=zref.z0_nm,
                box_z_nm=zref.box_z_nm,
                z_axis_direction=zref.direction,
                z_axis_source=zref.source,
                z_axis_surface_charge_uC_cm2=zref.surface_charge_uC_cm2,
            )
        )
    return payloads


def _validate_t0(payloads: list[CasePayload], out_dir: Path) -> dict[str, Any]:
    result: dict[str, Any] = {
        "available": False,
        "ok": False,
        "tolerance_rmsd_nm": 1.0e-6,
        "cases": [],
        "system_gro_md5": [],
        "trajectory_first_frame": {},
        "reason": None,
    }
    gro_hashes: list[tuple[CasePayload, str, Path]] = []
    for item in payloads:
        system_gro = item.relax_dir.parent / "02_system" / "system.gro"
        if not system_gro.is_file():
            result["cases"].append({"case": item.label, "available": False, "reason": "missing_02_system_system_gro", "path": str(system_gro)})
            continue
        digest = hashlib.md5(system_gro.read_bytes()).hexdigest()
        gro_hashes.append((item, digest, system_gro))
        result["system_gro_md5"].append({"case": item.label, "md5": digest, "path": str(system_gro)})
        result["cases"].append({"case": item.label, "available": True, "system_gro_md5": digest, "path": str(system_gro)})
    if len(gro_hashes) != len(payloads):
        result["reason"] = "missing_shared_t0_system_gro"
        _write_json(out_dir / "shared_t0_validation.json", result)
        return result
    ref_digest = gro_hashes[0][1]
    system_ok = all(digest == ref_digest for _item, digest, _path in gro_hashes)
    result["available"] = True
    result["ok"] = bool(system_ok)
    result["reason"] = None if system_ok else "02_system_system_gro_not_identical"

    # The exact shared-t0 contract is the case-level system.gro/topology.  The
    # first saved trajectory frame is only a diagnostic, because an xtc may start
    # after several MD steps depending on output cadence.
    try:
        import mdtraj as md
    except Exception as exc:
        result["trajectory_first_frame"] = {"available": False, "reason": f"mdtraj_unavailable: {exc}"}
        _write_json(out_dir / "shared_t0_validation.json", result)
        return result
    frames: list[tuple[CasePayload, np.ndarray, np.ndarray, float]] = []
    for item in payloads:
        base = _final_nvt_dir(item)
        xtc = base / "md.xtc"
        gro = base / "md.gro"
        if not xtc.is_file() or not gro.is_file():
            result.setdefault("trajectory_cases", []).append({"case": item.label, "available": False, "reason": "missing_final_nvt_xtc_or_gro"})
            continue
        try:
            frame = md.load_frame(str(xtc), 0, top=str(gro))
            frames.append((item, frame.xyz[0].copy(), frame.unitcell_lengths[0].copy(), float(frame.time[0])))
            result.setdefault("trajectory_cases", []).append(
                {
                    "case": item.label,
                    "available": True,
                    "natoms": int(frame.n_atoms),
                    "time_ps": float(frame.time[0]),
                    "box_nm": [float(x) for x in frame.unitcell_lengths[0][:3]],
                }
            )
        except Exception as exc:
            result.setdefault("trajectory_cases", []).append({"case": item.label, "available": False, "reason": str(exc)})
    if len(frames) != len(payloads):
        result["trajectory_first_frame"] = {"available": False, "reason": "missing_first_frames"}
        _write_json(out_dir / "shared_t0_validation.json", result)
        return result
    ref_item, ref_xyz, ref_box, _ref_time = frames[0]
    comparisons = []
    traj_ok = True
    for item, xyz, box, _time_ps in frames:
        if xyz.shape != ref_xyz.shape:
            rmsd = math.inf
            max_abs = math.inf
        else:
            delta = xyz - ref_xyz
            rmsd = float(np.sqrt(np.mean(delta * delta)))
            max_abs = float(np.max(np.abs(delta)))
        box_delta = [float(x) for x in (box - ref_box)]
        row = {
            "case": item.label,
            "reference_case": ref_item.label,
            "first_frame_rmsd_nm": rmsd,
            "max_abs_coord_delta_nm": max_abs,
            "box_delta_nm": box_delta,
            "ok": bool(rmsd <= result["tolerance_rmsd_nm"] and max(abs(x) for x in box_delta) <= 1.0e-6),
        }
        comparisons.append(row)
        traj_ok = traj_ok and bool(row["ok"])
    result["trajectory_first_frame"] = {
        "available": True,
        "ok": bool(traj_ok),
        "comparisons": comparisons,
        "note": "Diagnostic only; shared t=0 acceptance is based on identical 02_system/system.gro files.",
    }
    _write_json(out_dir / "shared_t0_validation.json", result)
    return result


def _write_json(path: Path, payload: dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")


def _aggregate_z_profiles(payloads: list[CasePayload], out_dir: Path) -> dict[str, dict[str, np.ndarray]]:
    by_case: dict[str, dict[str, dict[float, float]]] = {}
    all_x: list[float] = []
    for item in payloads:
        rows = _read_rows(item.analysis_dir / "z_density_profiles.csv")
        grouped: dict[str, dict[float, float]] = {}
        out_rows = []
        for row in rows:
            z_abs = _safe_float(row.get("z_mid_nm"))
            if not np.isfinite(z_abs):
                continue
            z_rel = _z_plot_nm(item, z_abs)
            species = _canonical_species(row.get("entity"))
            if str(row.get("entity_kind", "")).lower() == "phase":
                species = _canonical_species(row.get("entity"))
            if species not in SPECIES_ORDER:
                continue
            density = _safe_float(row.get("mass_density_g_cm3"), 0.0)
            grouped.setdefault(species, {})[round(z_rel, 6)] = grouped.setdefault(species, {}).get(round(z_rel, 6), 0.0) + float(density)
            out_rows.append(
                {
                    **row,
                    "z_abs_nm": z_abs,
                    "z_rel_nm": z_rel,
                    "z_plot_nm": z_rel,
                    "z_origin_nm": item.z0_nm,
                    "z_axis_direction": item.z_axis_direction,
                    "box_z_nm": item.box_z_nm,
                    "canonical_species": species,
                }
            )
            all_x.append(z_rel)
        _write_rows(out_dir / "zrel_csv" / f"{item.dirname}_z_density_profiles_zrel.csv", out_rows)
        by_case[item.label] = grouped
    arrays: dict[str, dict[str, np.ndarray]] = {}
    for item in payloads:
        arrays[item.label] = {}
        for species in SPECIES_ORDER:
            points = by_case.get(item.label, {}).get(species, {})
            if not points:
                arrays[item.label][species] = np.empty((0, 2), dtype=float)
                continue
            x = np.asarray(sorted(points), dtype=float)
            y = np.asarray([points[float(round(v, 6))] for v in x], dtype=float)
            arrays[item.label][species] = np.column_stack([x, y])
    return arrays


def _common_limits(series: Iterable[np.ndarray], *, y_pad: float = 0.08) -> tuple[tuple[float, float], tuple[float, float]]:
    xs: list[float] = []
    ys: list[float] = []
    for arr in series:
        if arr.size == 0:
            continue
        xs.extend([float(np.nanmin(arr[:, 0])), float(np.nanmax(arr[:, 0]))])
        ys.extend([float(np.nanmin(arr[:, 1])), float(np.nanmax(arr[:, 1]))])
    if not xs:
        return (0.0, 1.0), (0.0, 1.0)
    xlim = (float(min(xs)), float(max(xs)))
    ymin, ymax = float(min(ys)), float(max(ys))
    if ymax <= ymin:
        ymax = ymin + 1.0
    pad = (ymax - ymin) * y_pad
    return xlim, (max(0.0, ymin - pad), ymax + pad)


def _adaptive_fraction_ylim(series: Iterable[np.ndarray], *, floor_upper: float = 0.05) -> tuple[float, float]:
    """Return a readable y limit for membrane-fraction panels.

    The membrane fraction is a number fraction.  Keeping every panel at 0-1 is
    formally comparable, but it hides real sub-percent changes.  This helper
    keeps ranges comparable within a plotted group while making small uptake
    differences visible.
    """

    values: list[float] = []
    for arr in series:
        if arr.size == 0:
            continue
        col = np.asarray(arr[:, 1], dtype=float)
        values.extend([float(v) for v in col[np.isfinite(col)]])
    if not values:
        return 0.0, float(floor_upper)
    ymin, ymax = min(values), max(values)
    if ymax <= 1.0e-12:
        return 0.0, float(floor_upper)
    span = max(ymax - ymin, ymax * 0.20, floor_upper * 0.20)
    lo = max(0.0, ymin - 0.10 * span)
    hi = min(1.0, max(float(floor_upper), ymax + 0.20 * span))
    if hi <= lo:
        hi = min(1.0, lo + float(floor_upper))
    return float(lo), float(hi)


def _plot_z_facets_by_charge(z_arrays: dict[str, dict[str, np.ndarray]], payloads: list[CasePayload], fig_dir: Path) -> dict[str, Path]:
    _configure_matplotlib()
    import matplotlib.pyplot as plt

    fig, axes = plt.subplots(2, 2, figsize=(11.5, 6.6), sharex=True, sharey=True)
    all_series = [z_arrays[item.label][sp] for item in payloads for sp in SPECIES_ORDER if sp != "CMCNA"]
    xlim, ylim = _common_limits(all_series)
    for ax, item in zip(axes.flat, payloads):
        for sp in SPECIES_ORDER:
            arr = z_arrays[item.label].get(sp, np.empty((0, 2)))
            if arr.size == 0:
                continue
            ax.fill_between(arr[:, 0], arr[:, 1], step="mid", alpha=0.18, color=_species_color(sp))
            ax.step(arr[:, 0], arr[:, 1], where="mid", lw=1.2, color=_species_color(sp), label=sp)
        ax.set_title(item.label)
        ax.set_xlim(*xlim)
        ax.set_ylim(*ylim)
        ax.grid(True, alpha=0.20)
    axes[0, 0].legend(ncol=4, frameon=False, fontsize=7)
    fig.supxlabel("z_plot / nm (0 = negative CMC-facing graphite; +z enters CMC)")
    fig.supylabel("mass density / g cm$^{-3}$")
    fig.suptitle("z distribution by charge: species compared within each charge state")
    fig.tight_layout()
    out = _save_fig(fig, fig_dir / "z_distribution_by_charge_facets")
    plt.close(fig)
    return out


def _plot_z_facets_by_charge_zoom(
    z_arrays: dict[str, dict[str, np.ndarray]],
    payloads: list[CasePayload],
    fig_dir: Path,
    *,
    zoom_nm: float | None = None,
) -> dict[str, Path]:
    _configure_matplotlib()
    import matplotlib.pyplot as plt

    zoom_nm = float(zoom_nm if zoom_nm is not None else os.environ.get("EG08_ZOOM_CMC_INTERFACE_NM", "3.0"))
    zoom_nm = max(0.2, zoom_nm)
    fig, axes = plt.subplots(2, 2, figsize=(11.5, 6.6), sharex=True, sharey=True)
    rows_out: list[dict[str, Any]] = []
    for ax, item in zip(axes.flat, payloads):
        for sp in SPECIES_ORDER:
            arr = z_arrays[item.label].get(sp, np.empty((0, 2)))
            if arr.size == 0:
                continue
            mask = np.isfinite(arr[:, 0]) & np.isfinite(arr[:, 1]) & (arr[:, 0] >= 0.0) & (arr[:, 0] <= zoom_nm)
            if not np.any(mask):
                continue
            sub = arr[mask]
            ax.fill_between(sub[:, 0], sub[:, 1], step="mid", alpha=0.16, color=_species_color(sp))
            ax.step(sub[:, 0], sub[:, 1], where="mid", lw=1.25, color=_species_color(sp), label=sp)
            for x, y in sub:
                rows_out.append({"case": item.label, "charge_uC_cm2": item.charge, "species": sp, "z_plot_nm": float(x), "mass_density_g_cm3": float(y)})
        ax.set_title(item.label)
        ax.set_xlim(0.0, zoom_nm)
        ax.set_ylim(0.0, 0.5)
        ax.grid(True, alpha=0.20)
    axes[0, 0].legend(ncol=4, frameon=False, fontsize=7)
    fig.supxlabel("z_plot / nm (CMC-side graphite surface at 0)")
    fig.supylabel("mass density / g cm$^{-3}$")
    fig.suptitle(f"CMC-interface z distribution zoom (0-{zoom_nm:g} nm, y=0-0.5)")
    fig.tight_layout()
    _write_rows(fig_dir / "z_distribution_by_charge_facets_cmc_interface_zoom.csv", rows_out)
    out = _save_fig(fig, fig_dir / "z_distribution_by_charge_facets_cmc_interface_zoom")
    plt.close(fig)
    return out


def _plot_z_facets_by_species(z_arrays: dict[str, dict[str, np.ndarray]], payloads: list[CasePayload], fig_dir: Path) -> dict[str, Path]:
    _configure_matplotlib()
    import matplotlib.pyplot as plt

    species = ("EC", "EMC", "DEC", "LI", "PF6", "NA")
    fig, axes = plt.subplots(3, 2, figsize=(11.0, 8.0), sharex=True)
    all_series = [z_arrays[item.label][sp] for item in payloads for sp in species]
    xlim, _ylim = _common_limits(all_series)
    ymax_by_species = {}
    for sp in species:
        _x, y = _common_limits([z_arrays[item.label].get(sp, np.empty((0, 2))) for item in payloads])
        ymax_by_species[sp] = y
    for ax, sp in zip(axes.flat, species):
        for item in payloads:
            arr = z_arrays[item.label].get(sp, np.empty((0, 2)))
            if arr.size == 0:
                continue
            ax.step(arr[:, 0], arr[:, 1], where="mid", lw=1.5, color=_charge_color(item.charge), label=item.label)
        ax.set_title(sp)
        ax.set_xlim(*xlim)
        ax.set_ylim(*ymax_by_species[sp])
        ax.grid(True, alpha=0.20)
    axes[0, 0].legend(ncol=2, frameon=False, fontsize=7)
    fig.supxlabel("z_plot / nm (0 = negative CMC-facing graphite; +z enters CMC)")
    fig.supylabel("mass density / g cm$^{-3}$")
    fig.suptitle("z distribution by species: charge sweep compared within each species")
    fig.tight_layout()
    out = _save_fig(fig, fig_dir / "z_distribution_by_species_facets")
    plt.close(fig)
    return out


def _membrane_timeseries(payloads: list[CasePayload], out_dir: Path) -> dict[str, dict[str, np.ndarray]]:
    data: dict[str, dict[str, list[tuple[float, float, float, float, float]]]] = {}
    for item in payloads:
        rows = _read_rows(item.analysis_dir / "membrane_permeation_timeseries.csv")
        by_bin: dict[tuple[str, int], list[tuple[float, float, float, float]]] = {}
        for row in rows:
            sp = _canonical_species(row.get("species") or row.get("moltype"))
            if sp not in PERMEANT_SPECIES and sp != "NA":
                continue
            t_ns = _safe_float(row.get("time_ps")) / 1000.0
            if not np.isfinite(t_ns):
                continue
            bin_idx = int(math.floor(t_ns / 1.0))
            feed = _safe_float(row.get("feed_count"), 0.0)
            mem = _safe_float(row.get("membrane_count"), 0.0)
            perm = _safe_float(row.get("permeate_count"), 0.0)
            entries = _safe_float(row.get("cumulative_entry_events"), 0.0)
            by_bin.setdefault((sp, bin_idx), []).append((feed, mem, perm, entries))
        for (sp, bin_idx), vals in by_bin.items():
            arr = np.asarray(vals, dtype=float)
            feed, mem, perm, entries = np.nanmean(arr, axis=0)
            total = feed + mem + perm
            frac = mem / total if total > 0 else np.nan
            data.setdefault(item.label, {}).setdefault(sp, []).append((bin_idx + 0.5, frac, entries, mem, total))
    final: dict[str, dict[str, np.ndarray]] = {}
    rows_out = []
    for label, by_sp in data.items():
        final[label] = {}
        for sp, vals in by_sp.items():
            vals = sorted(vals)
            arr = np.asarray(vals, dtype=float)
            if arr.shape[0] >= 3:
                smooth = np.convolve(arr[:, 1], np.ones(3) / 3.0, mode="same")
                arr[:, 1] = smooth
            final[label][sp] = arr
            for row in arr:
                rows_out.append(
                    {
                        "case": label,
                        "species": sp,
                        "time_ns": float(row[0]),
                        "membrane_fraction_smoothed": float(row[1]),
                        "cumulative_entry_events": float(row[2]),
                        "membrane_count": float(row[3]),
                        "total_count": float(row[4]),
                    }
                )
    _write_rows(out_dir / "membrane_fraction_1ns_3point_smooth.csv", rows_out)
    return final


def _plot_membrane_fraction(data: dict[str, dict[str, np.ndarray]], payloads: list[CasePayload], fig_dir: Path) -> tuple[dict[str, Path], dict[str, Path]]:
    _configure_matplotlib()
    import matplotlib.pyplot as plt

    by_species: dict[str, Path] = {}
    for sp in PERMEANT_SPECIES:
        fig, ax = plt.subplots(figsize=(7.0, 3.8))
        ylim = _adaptive_fraction_ylim([data.get(item.label, {}).get(sp, np.empty((0, 5))) for item in payloads])
        for item in payloads:
            arr = data.get(item.label, {}).get(sp, np.empty((0, 5)))
            if arr.size:
                ax.plot(arr[:, 0], arr[:, 1], lw=1.8, color=_charge_color(item.charge), label=item.label)
        ax.set_ylim(*ylim)
        ax.set_xlabel("time / ns")
        ax.set_ylabel(f"f_mem({sp}) = N_mem / N_total")
        ax.set_title(f"{sp}: membrane fraction across charge states")
        ax.grid(True, alpha=0.25)
        ax.legend(frameon=False, ncol=2)
        fig.tight_layout()
        out = _save_fig(fig, fig_dir / f"membrane_fraction_species_{sp}")
        by_species[sp] = out["png"]
        plt.close(fig)

    by_charge: dict[str, Path] = {}
    for item in payloads:
        fig, ax = plt.subplots(figsize=(7.0, 3.8))
        ylim = _adaptive_fraction_ylim([data.get(item.label, {}).get(sp, np.empty((0, 5))) for sp in PERMEANT_SPECIES])
        for sp in PERMEANT_SPECIES:
            arr = data.get(item.label, {}).get(sp, np.empty((0, 5)))
            if arr.size:
                ax.plot(arr[:, 0], arr[:, 1], lw=1.8, color=_species_color(sp), label=sp)
        ax.set_ylim(*ylim)
        ax.set_xlabel("time / ns")
        ax.set_ylabel("f_mem = N_mem / N_total")
        ax.set_title(f"{item.label}: membrane fraction by species")
        ax.grid(True, alpha=0.25)
        ax.legend(frameon=False, ncol=5)
        fig.tight_layout()
        out = _save_fig(fig, fig_dir / f"membrane_fraction_charge_{item.dirname}")
        by_charge[item.label] = out["png"]
        plt.close(fig)
    return by_species, by_charge


def _penetration_metrics(payloads: list[CasePayload], out_dir: Path) -> list[dict[str, Any]]:
    metrics: list[dict[str, Any]] = []
    for item in payloads:
        summary = _read_json(item.analysis_dir / "membrane_permeation_summary.json")
        by_species = dict(summary.get("summary_by_species") or {})
        hist_rows = _read_rows(item.analysis_dir / "penetration_depth_distribution.csv")
        hist_by_sp: dict[str, list[dict[str, str]]] = {}
        for row in hist_rows:
            sp = _canonical_species(row.get("species"))
            hist_by_sp.setdefault(sp, []).append(row)
        for sp in PERMEANT_SPECIES:
            rec = dict(by_species.get(sp) or {})
            rows = hist_by_sp.get(sp, [])
            depths = np.asarray([_safe_float(row.get("depth_mid_nm")) for row in rows], dtype=float)
            pct = np.asarray([_safe_float(row.get("percent_of_penetrated_frames")) for row in rows], dtype=float) / 100.0
            finite = np.isfinite(depths) & np.isfinite(pct)
            depths = depths[finite]
            pct = pct[finite]
            auc = _trapz(pct, depths) if depths.size else np.nan
            d95 = np.nan
            if depths.size and np.nansum(pct) > 0:
                order = np.argsort(depths)
                cum = np.cumsum(pct[order]) / np.nansum(pct)
                idx = min(int(np.searchsorted(cum, 0.95, side="left")), int(depths.size - 1))
                d95 = float(depths[order][idx])
            initial_feed = _safe_float(rec.get("initial_feed_count"), 0.0)
            entry_count = _safe_float(rec.get("entry_event_count"), 0.0)
            metrics.append(
                {
                    "case": item.label,
                    "charge_uC_cm2": item.charge,
                    "species": sp,
                    "P_entry": entry_count / initial_feed if initial_feed > 0 else np.nan,
                    "entry_event_count": entry_count,
                    "translocation_event_count": _safe_float(rec.get("translocation_event_count"), 0.0),
                    "D_mean_nm": _safe_float(rec.get("mean_membrane_depth_nm"), np.nan),
                    "D_max_nm": _safe_float(rec.get("max_membrane_depth_nm"), np.nan),
                    "D95_nm": d95,
                    "AUC_depth_nm": auc,
                    "loading_molecules_nm3": _safe_float(rec.get("mean_membrane_loading_molecules_nm3"), np.nan),
                    "apparent_entry_flux_events_nm2_ns": _safe_float(rec.get("apparent_entry_flux_events_nm2_ns"), np.nan),
                }
            )
    _write_rows(out_dir / "penetration_capability_metrics.csv", metrics)
    return metrics


def _plot_penetration_metrics(metrics: list[dict[str, Any]], payloads: list[CasePayload], fig_dir: Path) -> dict[str, Path]:
    _configure_matplotlib()
    import matplotlib.pyplot as plt

    outputs: dict[str, Path] = {}
    for metric, ylabel in (
        ("P_entry", "entry events / initial feed"),
        ("D95_nm", "D95 penetration depth / nm"),
        ("AUC_depth_nm", "depth-distribution AUC / nm"),
        ("loading_molecules_nm3", "mean loading / molecules nm$^{-3}$"),
    ):
        fig, axes = plt.subplots(2, 2, figsize=(10.5, 6.4), sharey=True)
        vals = [_safe_float(row.get(metric)) for row in metrics]
        finite = [v for v in vals if np.isfinite(v)]
        ymax = max(finite) * 1.18 if finite else 1.0
        for ax, item in zip(axes.flat, payloads):
            rows = [row for row in metrics if row["case"] == item.label]
            x = np.arange(len(PERMEANT_SPECIES), dtype=float)
            y = [_safe_float(next((row.get(metric) for row in rows if row["species"] == sp), np.nan), np.nan) for sp in PERMEANT_SPECIES]
            ax.bar(x, y, color=[_species_color(sp) for sp in PERMEANT_SPECIES], width=0.72)
            ax.set_xticks(x, PERMEANT_SPECIES, rotation=0)
            ax.set_ylim(0.0, ymax if ymax > 0 else 1.0)
            ax.set_title(item.label)
            ax.grid(True, axis="y", alpha=0.25)
        fig.supxlabel("species")
        fig.supylabel(ylabel)
        fig.suptitle(f"Penetration capability metric: {metric}")
        fig.tight_layout()
        out = _save_fig(fig, fig_dir / f"penetration_metric_{metric}")
        outputs[metric] = out["png"]
        plt.close(fig)
    return outputs


def _draw_penetration_schematic(fig_dir: Path) -> dict[str, Path]:
    _configure_matplotlib()
    import matplotlib.pyplot as plt
    from matplotlib.patches import FancyArrowPatch, Rectangle

    fig, ax = plt.subplots(figsize=(9.8, 3.5))
    ax.set_xlim(0.0, 10.0)
    ax.set_ylim(0.0, 3.6)
    ax.axis("off")
    regions = [
        (0.4, 0.95, 2.7, "#DCEBFA", "feed\n(electrolyte side)"),
        (3.1, 0.95, 3.5, "#E9E1C9", "CMC membrane"),
        (6.6, 0.95, 2.9, "#E8F3E8", "permeate side"),
    ]
    for x, y, w, color, label in regions:
        ax.add_patch(Rectangle((x, y), w, 1.55, facecolor=color, edgecolor="#555555", lw=1.0))
        ax.text(x + 0.5 * w, y + 0.78, label, ha="center", va="center", fontsize=11)
    ax.add_patch(FancyArrowPatch((1.2, 2.85), (3.55, 2.85), arrowstyle="->", mutation_scale=14, lw=1.7, color="#2B5C8A"))
    ax.text(2.1, 3.05, "entry event\nfeed -> membrane", ha="center", fontsize=9, color="#2B5C8A")
    ax.add_patch(FancyArrowPatch((1.1, 0.55), (8.8, 0.55), arrowstyle="->", mutation_scale=14, lw=1.7, color="#7F3C8D"))
    ax.text(5.0, 0.15, "translocation = feed -> membrane -> permeate", ha="center", fontsize=9, color="#7F3C8D")
    ax.add_patch(FancyArrowPatch((3.1, 2.55), (4.65, 2.55), arrowstyle="<->", mutation_scale=12, lw=1.4, color="#E76F51"))
    ax.text(3.88, 2.73, "penetration depth d", ha="center", fontsize=9, color="#E76F51")
    ax.text(
        0.55,
        3.35,
        "P_entry = entry_event_count / initial_feed_count;  D95 = 95th percentile of d;  AUC_depth = integral of normalized depth distribution",
        ha="left",
        va="center",
        fontsize=9,
    )
    fig.tight_layout()
    out = _save_fig(fig, fig_dir / "penetration_metric_schematic", dpi=220)
    plt.close(fig)
    return out


def _orientation_rows(payloads: list[CasePayload], fig_dir: Path) -> tuple[list[dict[str, Any]], list[dict[str, Any]], dict[str, Path]]:
    """Collect EC/EMC/DEC carbonyl orientation in graphite EDL by charge state."""

    _configure_matplotlib()
    import matplotlib.pyplot as plt

    solvents = ("EC", "EMC", "DEC")
    angle_rows: list[dict[str, Any]] = []
    summary_rows: list[dict[str, Any]] = []
    for item in payloads:
        rows = _read_rows(item.analysis_dir / "adsorbed_orientation.csv")
        for sp in solvents:
            selected = []
            for row in rows:
                if _canonical_species(row.get("moltype") or row.get("species")) != sp:
                    continue
                if str(row.get("adsorbed", "")).lower() not in {"true", "1", "yes"}:
                    continue
                if str(row.get("orientation_available", "")).lower() not in {"true", "1", "yes"}:
                    continue
                angle = _safe_float(row.get("carbonyl_angle_deg"))
                if np.isfinite(angle):
                    selected.append(angle)
                    angle_rows.append(
                        {
                            "case": item.label,
                            "charge_uC_cm2": item.charge,
                            "species": sp,
                            "carbonyl_angle_deg": float(angle),
                            "nearest_graphite_phase": row.get("nearest_graphite_phase"),
                            "nearest_graphite_side": row.get("nearest_graphite_side"),
                        }
                    )
            arr = np.asarray(selected, dtype=float)
            if arr.size:
                q25, median, q75 = np.percentile(arr, [25, 50, 75])
                summary_rows.append(
                    {
                        "case": item.label,
                        "charge_uC_cm2": item.charge,
                        "species": sp,
                        "sample_count": int(arr.size),
                        "mean_angle_deg": float(np.mean(arr)),
                        "median_angle_deg": float(median),
                        "q25_angle_deg": float(q25),
                        "q75_angle_deg": float(q75),
                        "low_sample": bool(arr.size < 30),
                    }
                )
            else:
                summary_rows.append(
                    {
                        "case": item.label,
                        "charge_uC_cm2": item.charge,
                        "species": sp,
                        "sample_count": 0,
                        "mean_angle_deg": np.nan,
                        "median_angle_deg": np.nan,
                        "q25_angle_deg": np.nan,
                        "q75_angle_deg": np.nan,
                        "low_sample": True,
                    }
                )
    _write_rows(fig_dir / "carbonyl_orientation_samples.csv", angle_rows)
    _write_rows(fig_dir / "carbonyl_orientation_summary.csv", summary_rows)

    outputs: dict[str, Path] = {}
    if angle_rows:
        edges = np.linspace(0.0, 180.0, 19)
        for sp in solvents:
            fig, axes = plt.subplots(2, 2, figsize=(10.5, 6.2), sharex=True, sharey=True)
            has_any = False
            for ax, item in zip(axes.flat, payloads):
                vals = np.asarray(
                    [row["carbonyl_angle_deg"] for row in angle_rows if row["case"] == item.label and row["species"] == sp],
                    dtype=float,
                )
                if vals.size:
                    has_any = True
                    hist, _ = np.histogram(vals, bins=edges)
                    pct = 100.0 * hist.astype(float) / max(1, int(np.sum(hist)))
                    mid = 0.5 * (edges[:-1] + edges[1:])
                    ax.bar(mid, pct, width=np.diff(edges) * 0.88, color=_charge_color(item.charge), alpha=0.86)
                    ax.text(0.03, 0.92, f"n={vals.size}", transform=ax.transAxes, fontsize=8)
                else:
                    ax.text(0.50, 0.50, "insufficient samples", transform=ax.transAxes, ha="center", va="center", fontsize=9)
                ax.set_title(item.label)
                ax.set_xlim(0.0, 180.0)
                ax.grid(True, axis="y", alpha=0.22)
            if has_any:
                fig.supxlabel("C=O angle to graphite surface normal / deg")
                fig.supylabel("adsorbed oriented-frame fraction / %")
                fig.suptitle(f"{sp} carbonyl orientation in graphite EDL")
                fig.tight_layout()
                out = _save_fig(fig, fig_dir / f"carbonyl_orientation_distribution_{sp}")
                outputs[f"distribution_{sp}"] = out["png"]
            plt.close(fig)

        fig, axes = plt.subplots(1, 3, figsize=(12.0, 3.8), sharey=True)
        for ax, sp in zip(axes, solvents):
            xs: list[float] = []
            means: list[float] = []
            errs_low: list[float] = []
            errs_high: list[float] = []
            colors: list[str] = []
            for item in payloads:
                row = next((r for r in summary_rows if r["case"] == item.label and r["species"] == sp), None)
                if not row or not np.isfinite(_safe_float(row.get("mean_angle_deg"))):
                    continue
                xs.append(float(item.charge))
                mean = _safe_float(row["mean_angle_deg"])
                means.append(mean)
                errs_low.append(max(0.0, mean - _safe_float(row.get("q25_angle_deg"), mean)))
                errs_high.append(max(0.0, _safe_float(row.get("q75_angle_deg"), mean) - mean))
                colors.append(_charge_color(item.charge))
            if xs:
                ax.errorbar(xs, means, yerr=[errs_low, errs_high], fmt="o-", lw=1.6, color="#303030", ecolor="#606060", capsize=3)
                ax.scatter(xs, means, c=colors, s=48, zorder=3)
            ax.set_title(sp)
            ax.set_xlabel("CMC-facing charge / uC cm$^{-2}$")
            ax.set_ylim(0.0, 180.0)
            ax.grid(True, axis="y", alpha=0.24)
        axes[0].set_ylabel("mean C=O angle / deg")
        fig.suptitle("Graphite-EDL carbonyl orientation trend")
        fig.tight_layout()
        out = _save_fig(fig, fig_dir / "carbonyl_orientation_charge_trend")
        outputs["trend"] = out["png"]
        plt.close(fig)
    return summary_rows, angle_rows, outputs


def _find_topology_file(item: CasePayload, final_dir: Path) -> Path | None:
    candidates = [
        final_dir / "system.top",
        final_dir / "topol.top",
        item.relax_dir / "system.top",
        item.relax_dir / "topol.top",
        item.relax_dir.parent / "system.top",
        item.relax_dir.parent / "topol.top",
    ]
    for path in candidates:
        if path.is_file():
            return path
    for root in (item.relax_dir, item.relax_dir.parent):
        if not root.is_dir():
            continue
        for name in ("system.top", "topol.top"):
            hits = sorted(root.rglob(name))
            if hits:
                return hits[0]
    return None


def _membrane_bounds(item: CasePayload) -> tuple[float, float, str] | None:
    summary = _read_json(item.analysis_dir / "membrane_permeation_summary.json")
    lo = summary.get("membrane_z_lo_nm")
    hi = summary.get("membrane_z_hi_nm")
    feed_side = str(summary.get("feed_side") or "").lower()
    if lo is None or hi is None:
        nested = item.summary.get("membrane_permeation") if isinstance(item.summary, dict) else {}
        if isinstance(nested, dict):
            lo = lo if lo is not None else nested.get("membrane_z_lo_nm")
            hi = hi if hi is not None else nested.get("membrane_z_hi_nm")
            feed_side = feed_side or str(nested.get("feed_side") or "").lower()
    try:
        lo_f = float(lo)
        hi_f = float(hi)
    except Exception:
        return None
    if not np.isfinite(lo_f) or not np.isfinite(hi_f) or hi_f <= lo_f:
        return None
    if feed_side not in {"below", "above"}:
        feed_side = "above"
    return lo_f, hi_f, feed_side


def _atom_is_oxygen_like(name: str, atype: str, charge: float) -> bool:
    lname = str(name or "").strip().lower()
    latype = str(atype or "").strip().lower()
    return lname.startswith("o") or latype.startswith("o") or float(charge) <= -0.20


def _atom_is_fluorine_like(name: str, atype: str, charge: float) -> bool:
    lname = str(name or "").strip().lower()
    latype = str(atype or "").strip().lower()
    return lname.startswith("f") or latype.startswith("f") or float(charge) < -0.10


def _li_solvation_indices(instances: list[dict[str, Any]]) -> dict[str, np.ndarray]:
    li: list[int] = []
    solvent_o: list[int] = []
    cmc_o: list[int] = []
    pf6_f: list[int] = []
    for inst in instances:
        moltype = str(inst.get("moltype") or "")
        kind = str(inst.get("kind") or "").lower()
        species = _canonical_species(moltype)
        if "polymer" in kind or "cmc" in moltype.lower():
            species = "CMCNA"
        idx = np.asarray(inst.get("atom_indices_0"), dtype=int)
        charges = np.asarray(inst.get("charges"), dtype=float)
        names = [str(x) for x in inst.get("atomnames") or []]
        atypes = [str(x) for x in inst.get("atomtypes") or []]
        if idx.size == 0:
            continue
        if species == "LI" or (idx.size == 1 and ("li" in moltype.lower() or (charges.size and float(charges[0]) > 0.3))):
            li.append(int(idx[0]))
            continue
        if species in {"EC", "EMC", "DEC"}:
            candidates: list[tuple[float, int]] = []
            for local, atom_idx in enumerate(idx):
                charge = float(charges[local]) if local < charges.size else 0.0
                name = names[local] if local < len(names) else ""
                atype = atypes[local] if local < len(atypes) else ""
                if _atom_is_oxygen_like(name, atype, charge):
                    candidates.append((charge, int(atom_idx)))
            if candidates:
                solvent_o.append(int(sorted(candidates, key=lambda item: item[0])[0][1]))
            continue
        if species == "CMCNA":
            for local, atom_idx in enumerate(idx):
                charge = float(charges[local]) if local < charges.size else 0.0
                name = names[local] if local < len(names) else ""
                atype = atypes[local] if local < len(atypes) else ""
                if _atom_is_oxygen_like(name, atype, charge):
                    cmc_o.append(int(atom_idx))
            continue
        if species == "PF6":
            for local, atom_idx in enumerate(idx):
                charge = float(charges[local]) if local < charges.size else 0.0
                name = names[local] if local < len(names) else ""
                atype = atypes[local] if local < len(atypes) else ""
                if _atom_is_fluorine_like(name, atype, charge):
                    pf6_f.append(int(atom_idx))
    return {
        "li": np.asarray(sorted(set(li)), dtype=int),
        "solvent_o": np.asarray(sorted(set(solvent_o)), dtype=int),
        "cmc_o": np.asarray(sorted(set(cmc_o)), dtype=int),
        "pf6_f": np.asarray(sorted(set(pf6_f)), dtype=int),
    }


def _count_xy_pbc_z_open_neighbors(
    *,
    coords: np.ndarray,
    center_index: int,
    target_indices: np.ndarray,
    box: tuple[float, float, float],
    cutoff_nm: float,
) -> int:
    if target_indices.size == 0:
        return 0
    target = np.asarray(coords[target_indices], dtype=float)
    center = np.asarray(coords[int(center_index)], dtype=float)
    delta = target - center[None, :]
    lx, ly = max(float(box[0]), 1.0e-12), max(float(box[1]), 1.0e-12)
    delta[:, 0] -= lx * np.round(delta[:, 0] / lx)
    delta[:, 1] -= ly * np.round(delta[:, 1] / ly)
    dist2 = np.einsum("ij,ij->i", delta, delta)
    return int(np.count_nonzero(dist2 <= float(cutoff_nm) ** 2))


def _li_solvation_depth_profiles(payloads: list[CasePayload], fig_dir: Path) -> tuple[list[dict[str, Any]], dict[str, Path]]:
    rows_out: list[dict[str, Any]] = []
    summary_rows: list[dict[str, Any]] = []
    cutoff_nm = float(os.environ.get("EG08_LI_SOLVATION_CUTOFF_NM", "0.32"))
    depth_bin_nm = float(os.environ.get("EG08_LI_SOLVATION_DEPTH_BIN_NM", "0.10"))
    frame_stride = max(1, int(os.environ.get("EG08_LI_SOLVATION_FRAME_STRIDE", os.environ.get("EG08_ANALYSIS_FRAME_STRIDE", "4"))))
    for item in payloads:
        bounds = _membrane_bounds(item)
        final_dir = _final_nvt_dir(item)
        top_path = _find_topology_file(item, final_dir)
        gro = final_dir / "md.gro"
        xtc = final_dir / "md.xtc"
        if bounds is None or top_path is None or not gro.is_file():
            summary_rows.append(
                {
                    "case": item.label,
                    "available": False,
                    "reason": "missing_membrane_bounds_topology_or_final_gro",
                    "final_dir": str(final_dir),
                }
            )
            continue
        try:
            from yadonpy.gmx.analysis.interface_profile import _atom_payload, _iter_frames
            from yadonpy.gmx.topology import parse_system_top

            top = parse_system_top(top_path)
            atom_payload = _atom_payload(top, top_path.parent)
            site_indices = _li_solvation_indices(list(atom_payload.get("instances") or []))
        except Exception as exc:
            summary_rows.append({"case": item.label, "available": False, "reason": f"topology_parse_failed: {exc}", "topology": str(top_path)})
            continue
        li_indices = site_indices["li"]
        if li_indices.size == 0:
            summary_rows.append({"case": item.label, "available": False, "reason": "no_li_indices", "topology": str(top_path)})
            continue
        membrane_lo, membrane_hi, feed_side = bounds
        thickness = max(0.0, membrane_hi - membrane_lo)
        edges = np.arange(0.0, thickness + depth_bin_nm * 1.5, depth_bin_nm, dtype=float)
        if edges.size < 2:
            edges = np.asarray([0.0, max(depth_bin_nm, thickness)], dtype=float)
        sums = {target: np.zeros(edges.size - 1, dtype=float) for target in LI_SOLVATION_TARGETS}
        sample_counts = np.zeros(edges.size - 1, dtype=float)
        frame_count = 0
        li_inside_total = 0
        try:
            frame_iter = _iter_frames(gro_path=gro, xtc_path=xtc if xtc.is_file() else None, frame_stride=frame_stride, chunk=25)
            for _time_ps, coords, box in frame_iter:
                frame_count += 1
                z = np.asarray(coords[li_indices, 2], dtype=float)
                inside = (z >= membrane_lo) & (z <= membrane_hi)
                if not np.any(inside):
                    continue
                for li_idx, z_li in zip(li_indices[inside], z[inside]):
                    depth = (membrane_hi - float(z_li)) if feed_side == "above" else (float(z_li) - membrane_lo)
                    if depth < 0.0:
                        continue
                    bin_idx = int(np.searchsorted(edges, depth, side="right") - 1)
                    if bin_idx < 0 or bin_idx >= sample_counts.size:
                        continue
                    sample_counts[bin_idx] += 1.0
                    li_inside_total += 1
                    for target in LI_SOLVATION_TARGETS:
                        sums[target][bin_idx] += float(
                            _count_xy_pbc_z_open_neighbors(
                                coords=coords,
                                center_index=int(li_idx),
                                target_indices=site_indices[target],
                                box=box,
                                cutoff_nm=cutoff_nm,
                            )
                        )
        except Exception as exc:
            summary_rows.append({"case": item.label, "available": False, "reason": f"trajectory_scan_failed: {exc}", "final_dir": str(final_dir)})
            continue
        for bin_idx in range(sample_counts.size):
            n = float(sample_counts[bin_idx])
            means = {target: (float(sums[target][bin_idx]) / n if n > 0 else np.nan) for target in LI_SOLVATION_TARGETS}
            total_cn = float(sum(v for v in means.values() if np.isfinite(v)))
            row = {
                "case": item.label,
                "charge_uC_cm2": float(item.charge),
                "depth_bin_lo_nm": float(edges[bin_idx]),
                "depth_bin_hi_nm": float(edges[bin_idx + 1]),
                "depth_mid_nm": float(0.5 * (edges[bin_idx] + edges[bin_idx + 1])),
                "li_frame_samples": int(n),
                "cutoff_nm": cutoff_nm,
                "feed_side": feed_side,
                "cn_solvent_o": means["solvent_o"],
                "cn_cmc_o": means["cmc_o"],
                "cn_pf6_f": means["pf6_f"],
                "cn_total": total_cn if n > 0 else np.nan,
                "fraction_solvent_o": means["solvent_o"] / total_cn if total_cn > 0 and np.isfinite(means["solvent_o"]) else np.nan,
                "fraction_cmc_o": means["cmc_o"] / total_cn if total_cn > 0 and np.isfinite(means["cmc_o"]) else np.nan,
                "fraction_pf6_f": means["pf6_f"] / total_cn if total_cn > 0 and np.isfinite(means["pf6_f"]) else np.nan,
            }
            rows_out.append(row)
        summary_rows.append(
            {
                "case": item.label,
                "available": True,
                "frame_count": int(frame_count),
                "frame_stride": int(frame_stride),
                "li_atom_count": int(li_indices.size),
                "li_inside_cmc_frame_samples": int(li_inside_total),
                "membrane_z_lo_nm": float(membrane_lo),
                "membrane_z_hi_nm": float(membrane_hi),
                "feed_side": feed_side,
                "cutoff_nm": cutoff_nm,
            }
        )
    _write_rows(fig_dir / "li_solvation_by_cmc_depth.csv", rows_out)
    _write_rows(fig_dir / "li_solvation_by_cmc_depth_summary.csv", summary_rows)
    plots = _plot_li_solvation_depth(rows_out, payloads, fig_dir)
    return summary_rows, plots


def _plot_li_solvation_depth(rows: list[dict[str, Any]], payloads: list[CasePayload], fig_dir: Path) -> dict[str, Path]:
    _configure_matplotlib()
    import matplotlib.pyplot as plt

    valid_rows = [row for row in rows if _safe_float(row.get("li_frame_samples"), 0.0) > 0]
    outputs: dict[str, Path] = {}
    if not valid_rows:
        return outputs
    cn_keys = {
        "solvent_o": "cn_solvent_o",
        "cmc_o": "cn_cmc_o",
        "pf6_f": "cn_pf6_f",
    }
    frac_keys = {
        "solvent_o": "fraction_solvent_o",
        "cmc_o": "fraction_cmc_o",
        "pf6_f": "fraction_pf6_f",
    }
    max_cn = max((_safe_float(row.get(key), 0.0) for row in valid_rows for key in cn_keys.values()), default=1.0)
    ymax = max(6.0, max_cn * 1.15)

    fig, axes = plt.subplots(2, 2, figsize=(11.0, 6.8), sharex=True, sharey=True)
    for ax, item in zip(axes.flat, payloads):
        case_rows = [row for row in valid_rows if row.get("case") == item.label]
        for target, key in cn_keys.items():
            x = np.asarray([_safe_float(row.get("depth_mid_nm")) for row in case_rows], dtype=float)
            y = np.asarray([_safe_float(row.get(key)) for row in case_rows], dtype=float)
            samples = np.asarray([_safe_float(row.get("li_frame_samples"), 0.0) for row in case_rows], dtype=float)
            mask = np.isfinite(x) & np.isfinite(y) & (samples > 0)
            if np.any(mask):
                ax.plot(x[mask], y[mask], marker="o", ms=3.0, lw=1.6, color=LI_SOLVATION_PALETTE[target], label=target.replace("_", "-"))
        ax.axhline(4.0, color="#808080", lw=0.9, ls=":", alpha=0.8)
        ax.axhline(6.0, color="#808080", lw=0.9, ls="--", alpha=0.8)
        ax.set_ylim(0.0, ymax)
        ax.set_title(item.label)
        ax.grid(True, alpha=0.22)
    axes[0, 0].legend(frameon=False, ncol=3, fontsize=7)
    fig.supxlabel("Li penetration depth into CMC from electrolyte-side boundary / nm")
    fig.supylabel("mean coordination count within cutoff")
    fig.suptitle("Depth-resolved Li solvation structure inside CMC membrane")
    fig.tight_layout()
    out = _save_fig(fig, fig_dir / "li_solvation_cn_by_cmc_depth")
    outputs["cn_depth"] = out["png"]
    plt.close(fig)

    fig, axes = plt.subplots(1, 3, figsize=(12.0, 3.7), sharex=True, sharey=True)
    for ax, target in zip(axes, LI_SOLVATION_TARGETS):
        key = cn_keys[target]
        for item in payloads:
            case_rows = [row for row in valid_rows if row.get("case") == item.label]
            x = np.asarray([_safe_float(row.get("depth_mid_nm")) for row in case_rows], dtype=float)
            y = np.asarray([_safe_float(row.get(key)) for row in case_rows], dtype=float)
            samples = np.asarray([_safe_float(row.get("li_frame_samples"), 0.0) for row in case_rows], dtype=float)
            mask = np.isfinite(x) & np.isfinite(y) & (samples > 0)
            if np.any(mask):
                ax.plot(x[mask], y[mask], marker="o", ms=2.8, lw=1.5, color=_charge_color(item.charge), label=item.label)
        ax.axhline(4.0, color="#808080", lw=0.9, ls=":", alpha=0.8)
        ax.axhline(6.0, color="#808080", lw=0.9, ls="--", alpha=0.8)
        ax.set_title(target.replace("_", "-"))
        ax.set_ylim(0.0, ymax)
        ax.grid(True, alpha=0.22)
    axes[0].legend(frameon=False, fontsize=7)
    fig.supxlabel("Li penetration depth into CMC from electrolyte-side boundary / nm")
    fig.supylabel("mean coordination count within cutoff")
    fig.suptitle("Depth-resolved Li solvation: charge sweep by coordinating site")
    fig.tight_layout()
    out = _save_fig(fig, fig_dir / "li_solvation_cn_by_site_charge_facets")
    outputs["cn_site_charge_facets"] = out["png"]
    plt.close(fig)

    fig, axes = plt.subplots(2, 2, figsize=(11.0, 6.8), sharex=True, sharey=True)
    for ax, item in zip(axes.flat, payloads):
        case_rows = [row for row in valid_rows if row.get("case") == item.label]
        for target, key in frac_keys.items():
            x = np.asarray([_safe_float(row.get("depth_mid_nm")) for row in case_rows], dtype=float)
            y = np.asarray([_safe_float(row.get(key)) for row in case_rows], dtype=float)
            samples = np.asarray([_safe_float(row.get("li_frame_samples"), 0.0) for row in case_rows], dtype=float)
            mask = np.isfinite(x) & np.isfinite(y) & (samples > 0)
            if np.any(mask):
                ax.plot(x[mask], y[mask], marker="o", ms=3.0, lw=1.6, color=LI_SOLVATION_PALETTE[target], label=target.replace("_", "-"))
        ax.set_ylim(0.0, 1.0)
        ax.set_title(item.label)
        ax.grid(True, alpha=0.22)
    axes[0, 0].legend(frameon=False, ncol=3, fontsize=7)
    fig.supxlabel("Li penetration depth into CMC from electrolyte-side boundary / nm")
    fig.supylabel("coordination composition fraction")
    fig.suptitle("Depth-resolved Li solvation composition inside CMC membrane")
    fig.tight_layout()
    out = _save_fig(fig, fig_dir / "li_solvation_fraction_by_cmc_depth")
    outputs["fraction_depth"] = out["png"]
    plt.close(fig)
    return outputs


def _recompute_edl_on_plot_axis(rows: list[dict[str, str]], item: CasePayload) -> list[dict[str, Any]]:
    points: list[tuple[float, float, dict[str, str]]] = []
    for row in rows:
        z_abs = _safe_float(row.get("z_nm"))
        rho = _safe_float(row.get("charge_density_e_nm3"), 0.0)
        if np.isfinite(z_abs):
            points.append((_z_plot_nm(item, z_abs), rho, row))
    points.sort(key=lambda item_: item_[0])
    if not points:
        return []
    z = np.asarray([p[0] for p in points], dtype=float)
    rho = np.asarray([p[1] for p in points], dtype=float)
    dz_values = np.diff(z)
    dz = float(np.nanmedian(dz_values[dz_values > 0])) if np.any(dz_values > 0) else 0.05
    if not np.isfinite(dz) or dz <= 0:
        dz = 0.05
    integrated = np.cumsum(rho * dz)
    elementary_charge = 1.602176634e-19
    eps0 = 8.8541878128e-12
    electric_field = integrated * elementary_charge / 1.0e-18 / eps0
    phi_surface = -np.cumsum(electric_field * dz * 1.0e-9)
    if phi_surface.size:
        phi_surface -= phi_surface[0]
    ref_hi = min(1.5, max(0.8 + 2.0 * dz, 0.45 * item.box_z_nm))
    ref_mask = (z >= 0.8) & (z <= ref_hi)
    ref_label = "local_interior_0.8_1.5_nm"
    if np.count_nonzero(ref_mask) < 3:
        ref_mask = (z >= 0.20 * item.box_z_nm) & (z <= 0.45 * item.box_z_nm)
        ref_label = "periodic_interior_0.20_0.45_Lz"
    if np.count_nonzero(ref_mask) < 3:
        ref_mask = np.isfinite(phi_surface)
        ref_label = "period_mean_fallback"
    ref_value = float(np.nanmean(phi_surface[ref_mask])) if np.count_nonzero(ref_mask) else 0.0
    phi = phi_surface - ref_value

    out: list[dict[str, Any]] = []
    if z.size and z[0] > 1.0e-9:
        out.append(
            {
                "case": item.label,
                "charge_uC_cm2": item.charge,
                "z_abs_nm": "",
                "z_rel_nm": 0.0,
                "z_plot_nm": 0.0,
                "z_origin_nm": item.z0_nm,
                "z_axis_direction": item.z_axis_direction,
                "box_z_nm": item.box_z_nm,
                "charge_density_e_nm3": float(rho[0]),
                "integrated_charge_e_nm2": 0.0,
                "electric_field_V_m": 0.0,
                "electrostatic_potential_V": float(-ref_value),
                "electrostatic_potential_surface_ref_V": 0.0,
                "potential_reference_V": ref_value,
                "potential_reference_region": ref_label,
            }
        )
    for idx, (_z_plot, _rho, raw) in enumerate(points):
        z_abs = _safe_float(raw.get("z_nm"))
        out.append(
            {
                **raw,
                "case": item.label,
                "charge_uC_cm2": item.charge,
                "z_abs_nm": z_abs,
                "z_rel_nm": float(z[idx]),
                "z_plot_nm": float(z[idx]),
                "z_origin_nm": item.z0_nm,
                "z_axis_direction": item.z_axis_direction,
                "box_z_nm": item.box_z_nm,
                "charge_density_e_nm3": float(rho[idx]),
                "integrated_charge_e_nm2": float(integrated[idx]),
                "electric_field_V_m": float(electric_field[idx]),
                "electrostatic_potential_V": float(phi[idx]),
                "electrostatic_potential_surface_ref_V": float(phi_surface[idx]),
                "potential_reference_V": ref_value,
                "potential_reference_region": ref_label,
            }
        )
    return out


def _plot_edl_overlays(payloads: list[CasePayload], fig_dir: Path) -> dict[str, Path]:
    _configure_matplotlib()
    import matplotlib.pyplot as plt

    columns = (
        ("charge_density_e_nm3", "charge density / e nm$^{-3}$", "edl_charge_density_zrel"),
        ("integrated_charge_e_nm2", "integrated charge / e nm$^{-2}$", "edl_integrated_charge_zrel"),
        ("electrostatic_potential_V", "potential / V", "edl_potential_zrel"),
    )
    outputs: dict[str, Path] = {}
    series_by_col: dict[str, dict[str, tuple[np.ndarray, np.ndarray]]] = {col: {} for col, _yl, _name in columns}
    all_corrected_rows: list[dict[str, Any]] = []
    for item in payloads:
        rows = _read_rows(item.analysis_dir / "electrostatic_potential.csv")
        out_rows = _recompute_edl_on_plot_axis(rows, item)
        all_corrected_rows.extend(out_rows)
        _write_rows(fig_dir / "zrel_csv" / f"{item.dirname}_electrostatic_potential_zrel.csv", out_rows)
        for col, _yl, _name in columns:
            x = np.asarray([_safe_float(row.get("z_plot_nm")) for row in out_rows], dtype=float)
            y = np.asarray([_safe_float(row.get(col)) for row in out_rows], dtype=float)
            mask = np.isfinite(x) & np.isfinite(y)
            series_by_col[col][item.label] = (x[mask], y[mask])
    _write_rows(fig_dir / "edl_profiles_zaxis_corrected.csv", all_corrected_rows)
    xlim = (0.0, max((item.box_z_nm for item in payloads), default=1.0))
    for col, ylabel, name in columns:
        fig, ax = plt.subplots(figsize=(7.2, 4.0))
        ys = []
        for item in payloads:
            x, y = series_by_col[col].get(item.label, (np.asarray([]), np.asarray([])))
            if x.size:
                ax.plot(x, y, lw=1.8, color=_charge_color(item.charge), label=item.label)
                ys.extend([float(np.nanmin(y)), float(np.nanmax(y))])
        ax.set_xlim(*xlim)
        if ys:
            pad = (max(ys) - min(ys)) * 0.08 or 1.0
            ax.set_ylim(min(ys) - pad, max(ys) + pad)
        ax.set_xlabel("z_plot / nm (0 = negative CMC-facing graphite; +z enters CMC)")
        ax.set_ylabel(ylabel)
        ax.set_title(name.replace("_", " "))
        ax.grid(True, alpha=0.25)
        ax.legend(frameon=False, ncol=2)
        fig.tight_layout()
        out = _save_fig(fig, fig_dir / name)
        outputs[name] = out["png"]
        plt.close(fig)
    return outputs


def _plot_rdf_cn_last_window(payloads: list[CasePayload], fig_dir: Path) -> dict[str, Path]:
    _configure_matplotlib()
    import matplotlib.pyplot as plt

    pairs = ("cation-polymer_o", "cation-solvent_o", "cation-anion_f")
    outputs: dict[str, Path] = {}
    for pair in pairs:
        fig, ax = plt.subplots(figsize=(7.2, 4.1))
        ax_cn = ax.twinx()
        for item in payloads:
            rows = [row for row in _read_rows(item.analysis_dir / "time_series" / "rdf_cn_curves_timeseries.csv") if row.get("pair") == pair]
            if not rows:
                continue
            max_end = max(_safe_float(row.get("time_end_ps"), -1.0) for row in rows)
            selected = [row for row in rows if abs(_safe_float(row.get("time_end_ps"), -1.0) - max_end) < 1.0e-6]
            r = np.asarray([_safe_float(row.get("r_nm")) for row in selected], dtype=float)
            g = np.asarray([_safe_float(row.get("g_r")) for row in selected], dtype=float)
            cn = np.asarray([_safe_float(row.get("cn_r")) for row in selected], dtype=float)
            mask = np.isfinite(r) & np.isfinite(g) & np.isfinite(cn)
            if not np.any(mask):
                continue
            color = _charge_color(item.charge)
            ax.plot(r[mask], g[mask], lw=1.6, color=color, label=f"{item.label} RDF")
            ax_cn.plot(r[mask], cn[mask], lw=1.4, ls="--", color=color, alpha=0.9)
            if np.any(mask):
                peak_idx = int(np.nanargmax(g[mask]))
                rr = r[mask][peak_idx]
                gg = g[mask][peak_idx]
                ax.annotate(f"{rr:.2f} nm", xy=(rr, gg), xytext=(rr, gg * 1.05), fontsize=7, color=color)
        ax.set_xlabel("r / nm")
        ax.set_ylabel("RDF g(r), solid")
        ax_cn.set_ylabel("CN(r), dashed")
        ax_cn.set_ylim(0.0, 6.0)
        ax.set_title(f"Interface RDF/CN last window: {pair}")
        ax.grid(True, alpha=0.25)
        handles, labels = ax.get_legend_handles_labels()
        if handles:
            ax.legend(handles, labels, frameon=False, fontsize=7, ncol=2)
        else:
            ax.text(0.5, 0.5, "No RDF/CN rows available for this pair", transform=ax.transAxes, ha="center", va="center", fontsize=10)
        fig.tight_layout()
        out = _save_fig(fig, fig_dir / f"rdf_cn_last_window_{pair.replace('-', '_')}")
        outputs[pair] = out["png"]
        plt.close(fig)
    return outputs


def _add_textbox(slide, left, top, width, height, text: str, font_size: int = 13) -> None:
    from pptx.util import Pt

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, line in enumerate(str(text).splitlines() or [""]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        for run in p.runs:
            run.font.name = "Arial"
            run.font.size = Pt(font_size)


def _add_picture(slide, path: Path | None, left, top, width=None, height=None) -> bool:
    if path is None or not path.is_file():
        return False
    try:
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)
        return True
    except Exception:
        return False


def _add_movie_tile(slide, mp4_path: Path, poster: Path | None, left, top, width, height) -> None:
    from pptx.util import Inches

    shown = _add_picture(slide, poster, left, top, width=width, height=height)
    if mp4_path.is_file():
        try:
            slide.shapes.add_movie(
                str(mp4_path),
                left + width - Inches(0.82),
                top + Inches(0.07),
                Inches(0.72),
                Inches(0.50),
                poster_frame_image=str(poster) if poster and poster.is_file() else None,
                mime_type="video/mp4",
            )
        except Exception:
            pass
        _add_textbox(slide, left + Inches(0.07), top + height - Inches(0.28), width - Inches(0.14), Inches(0.22), f"MP4: {mp4_path.name}", font_size=7)
    elif not shown:
        _add_textbox(slide, left, top, width, height, f"Missing MP4/poster:\n{mp4_path}", font_size=9)


def _add_table(slide, rows: list[list[Any]], left, top, width, height, *, font_size: int = 8) -> None:
    from pptx.util import Pt

    if not rows:
        return
    table = slide.shapes.add_table(len(rows), len(rows[0]), left, top, width, height).table
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "Arial"
                    run.font.size = Pt(font_size)


def _result_slide(prs, title: str, image: Path | None, text: str, *, source_svg: Path | None = None):
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    shown = _add_picture(slide, image, Inches(0.45), Inches(0.82), width=Inches(7.25))
    if not shown:
        _add_table(
            slide,
            [["figure status", "not generated"], ["expected source", source_svg or image or "not available"]],
            Inches(0.75),
            Inches(1.25),
            Inches(6.4),
            Inches(1.1),
            font_size=9,
        )
    note = text
    if source_svg is not None:
        note += f"\n\nSVG source:\n{source_svg}"
    _add_textbox(slide, Inches(8.0), Inches(0.9), Inches(5.0), Inches(5.95), note, font_size=12)
    return slide


def _audit_ppt(path: Path, out_json: Path) -> dict[str, Any]:
    from pptx import Presentation

    prs = Presentation(str(path))
    slides = []
    ok = True
    for idx, slide in enumerate(prs.slides, start=1):
        title = ""
        pics = movies = charts = tables = 0
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and not title:
                title = shape.text_frame.text.strip().split("\n")[0]
            if shape.shape_type == 13:
                pics += 1
            elif shape.shape_type in (16, 17):
                movies += 1
            elif shape.shape_type == 3:
                charts += 1
            elif shape.shape_type == 19:
                tables += 1
        visual = pics + movies + charts + tables
        data_slide = idx > 2 and "Notes" not in title
        if data_slide and visual == 0:
            ok = False
        slides.append({"index": idx, "title": title, "pictures": pics, "movies": movies, "charts": charts, "tables": tables, "visual_count": visual})
    payload = {"ok": ok, "ppt": str(path), "slides": slides, "size_mb": path.stat().st_size / 1024 / 1024}
    _write_json(out_json, payload)
    return payload


def _audit_z_axis(payloads: list[CasePayload], fig_dir: Path, out_json: Path) -> dict[str, Any]:
    checks: list[dict[str, Any]] = []
    ok = True
    for item in payloads:
        case_checks = []
        for path in (
            fig_dir / "zrel_csv" / f"{item.dirname}_z_density_profiles_zrel.csv",
            fig_dir / "zrel_csv" / f"{item.dirname}_electrostatic_potential_zrel.csv",
        ):
            rows = _read_rows(path)
            values = np.asarray([_safe_float(row.get("z_plot_nm", row.get("z_rel_nm"))) for row in rows], dtype=float)
            values = values[np.isfinite(values)]
            file_ok = bool(values.size == 0 or (np.nanmin(values) >= -1.0e-9 and np.nanmax(values) <= item.box_z_nm + 1.0e-6))
            ok = ok and file_ok
            case_checks.append(
                {
                    "file": str(path),
                    "row_count": int(values.size),
                    "min_z_plot_nm": float(np.nanmin(values)) if values.size else None,
                    "max_z_plot_nm": float(np.nanmax(values)) if values.size else None,
                    "box_z_nm": item.box_z_nm,
                    "ok": file_ok,
                }
            )
        if item.charge < 0 and item.z_axis_surface_charge_uC_cm2 is not None:
            charge_ok = item.z_axis_surface_charge_uC_cm2 < 0
        else:
            charge_ok = True
        ok = ok and charge_ok
        checks.append(
            {
                "case": item.label,
                "z0_nm": item.z0_nm,
                "direction": item.z_axis_direction,
                "surface_charge_uC_cm2_at_z0": item.z_axis_surface_charge_uC_cm2,
                "negative_surface_charge_ok": charge_ok,
                "files": case_checks,
            }
        )
    payload = {
        "ok": bool(ok),
        "definition": "z_plot_nm is constrained to [0, Lz); x=0 is the CMC-facing graphite surface and +x points into CMC.",
        "cases": checks,
    }
    _write_json(out_json, payload)
    return payload


def _read_gro_coordinates(path: Path) -> tuple[list[str], np.ndarray, tuple[float, float, float]]:
    if not path.is_file():
        return [], np.empty((0, 3), dtype=float), (1.0, 1.0, 1.0)
    lines = path.read_text(encoding="utf-8", errors="replace").splitlines()
    if len(lines) < 3:
        return [], np.empty((0, 3), dtype=float), (1.0, 1.0, 1.0)
    try:
        n = int(lines[1].strip())
    except Exception:
        n = max(0, len(lines) - 3)
    species: list[str] = []
    coords: list[tuple[float, float, float]] = []
    for line in lines[2 : 2 + n]:
        resname = line[5:10].strip() if len(line) >= 10 else ""
        atomname = line[10:15].strip() if len(line) >= 15 else ""
        species.append(resname or atomname)
        try:
            coords.append((float(line[20:28]), float(line[28:36]), float(line[36:44])))
        except Exception:
            parts = line.split()
            if len(parts) >= 3:
                coords.append(tuple(float(x) for x in parts[-3:]))
    box = _box_from_gro(path)
    return species, np.asarray(coords, dtype=float), box


def _plot_structure_snapshot(path: Path, item: CasePayload, out_png: Path, *, title: str) -> Path | None:
    _configure_matplotlib()
    import matplotlib.pyplot as plt

    species, coords, box = _read_gro_coordinates(path)
    if coords.size == 0:
        return None
    z_plot = np.asarray([_z_plot_nm(item, z) for z in coords[:, 2]], dtype=float)
    y = coords[:, 1] % max(box[1], 1.0e-12)
    fig, ax = plt.subplots(figsize=(7.8, 3.8))
    draw_order = ("GRAPHITE", "CMCNA", "EC", "EMC", "DEC", "LI", "PF6", "NA")
    for sp in draw_order:
        if sp == "GRAPHITE":
            mask = np.asarray(["GRAPH" in s or "GR" == s for s in species], dtype=bool)
            color = "#333333"
            size = 2.2
            alpha = 0.25
        else:
            mask = np.asarray([_canonical_species(s) == sp for s in species], dtype=bool)
            color = _species_color(sp)
            size = 4.0 if sp in {"LI", "NA"} else 2.4
            alpha = 0.55
        if np.any(mask):
            ax.scatter(z_plot[mask], y[mask], s=size, c=color, alpha=alpha, linewidths=0, label=sp)
    ax.set_xlim(0.0, item.box_z_nm)
    ax.set_ylim(0.0, max(box[1], 1.0))
    ax.set_xlabel("z_plot / nm (CMC-facing graphite surface at 0)")
    ax.set_ylabel("wrapped y / nm")
    ax.set_title(title)
    ax.grid(True, alpha=0.18)
    ax.legend(loc="upper right", ncol=4, fontsize=6, frameon=False)
    fig.tight_layout()
    out_png.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_png, dpi=220, bbox_inches="tight")
    plt.close(fig)
    return out_png


def _plot_structure_frame(
    *,
    species: list[str],
    coords: np.ndarray,
    box: tuple[float, float, float],
    item: CasePayload,
    out_png: Path,
    title: str,
    time_ps: float | None = None,
    max_points: int = 45000,
) -> Path | None:
    _configure_matplotlib()
    import matplotlib.pyplot as plt

    if coords.size == 0:
        return None
    z_plot = np.asarray([_z_plot_nm(item, z) for z in coords[:, 2]], dtype=float)
    y = np.asarray(coords[:, 1], dtype=float) % max(float(box[1]), 1.0e-12)
    fig, ax = plt.subplots(figsize=(7.8, 3.8))
    draw_order = ("GRAPHITE", "CMCNA", "EC", "EMC", "DEC", "LI", "PF6", "NA")
    for sp in draw_order:
        if sp == "GRAPHITE":
            mask = np.asarray(["GRAPH" in s or "GR" == s for s in species], dtype=bool)
            color = "#333333"
            size = 1.8
            alpha = 0.22
        else:
            mask = np.asarray([_canonical_species(s) == sp for s in species], dtype=bool)
            color = _species_color(sp)
            size = 3.0 if sp in {"LI", "NA"} else 1.8
            alpha = 0.50
        indices = np.flatnonzero(mask)
        if indices.size == 0:
            continue
        if indices.size > max_points // 4:
            keep = np.linspace(0, indices.size - 1, max(1, max_points // 4), dtype=int)
            indices = indices[keep]
        ax.scatter(z_plot[indices], y[indices], s=size, c=color, alpha=alpha, linewidths=0, label=sp)
    subtitle = "" if time_ps is None or not np.isfinite(time_ps) else f" | t={time_ps / 1000.0:.2f} ns"
    ax.set_xlim(0.0, item.box_z_nm)
    ax.set_ylim(0.0, max(float(box[1]), 1.0))
    ax.set_xlabel("z_plot / nm")
    ax.set_ylabel("wrapped y / nm")
    ax.set_title(f"{title}{subtitle}\nLx={box[0]:.2f} nm, Ly={box[1]:.2f} nm, Lz={box[2]:.2f} nm")
    ax.grid(True, alpha=0.18)
    ax.legend(loc="upper right", ncol=4, fontsize=6, frameon=False)
    fig.tight_layout()
    out_png.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_png, dpi=170, bbox_inches="tight")
    plt.close(fig)
    return out_png


def _structure_assets(payloads: list[CasePayload], fig_dir: Path) -> dict[str, Path]:
    assets: dict[str, Path] = {}
    if not payloads:
        return assets
    first = payloads[0]
    shared_gro = first.relax_dir.parent / "02_system" / "system.gro"
    shared_png = fig_dir / "structure_overview" / "shared_t0_structure.png"
    out = _plot_structure_snapshot(shared_gro, first, shared_png, title="Shared t=0 structure projection")
    if out is not None:
        assets["shared_t0"] = out
    for item in payloads:
        final_gro = _final_nvt_dir(item) / "md.gro"
        final_png = fig_dir / "structure_overview" / f"{item.dirname}_final_structure.png"
        out = _plot_structure_snapshot(final_gro, item, final_png, title=f"{item.label} final structure projection")
        if out is not None:
            assets[f"final_{item.dirname}"] = out
    return assets


def _trajectory_assets(item: CasePayload, fig_dir: Path) -> tuple[Path | None, Path | None]:
    movie_dir = fig_dir / "trajectory_overview"
    mp4 = movie_dir / f"{item.dirname}_trajectory_overview.mp4"
    poster = movie_dir / item.dirname / "frame_000.png"
    if mp4.is_file() and poster.is_file():
        return mp4, poster
    final_dir = _final_nvt_dir(item)
    xtc = final_dir / "md.xtc"
    gro = final_dir / "md.gro"
    if not gro.is_file():
        return None, None
    species, gro_coords, gro_box = _read_gro_coordinates(gro)
    if not species or gro_coords.size == 0:
        return None, None
    case_dir = movie_dir / item.dirname
    case_dir.mkdir(parents=True, exist_ok=True)
    max_frames = max(3, int(os.environ.get("EG08_TRAJECTORY_OVERVIEW_MAX_FRAMES", "12")))
    candidate_stride = max(1, int(os.environ.get("EG08_TRAJECTORY_OVERVIEW_FRAME_STRIDE", "100")))
    max_candidates = max(max_frames, int(os.environ.get("EG08_TRAJECTORY_OVERVIEW_MAX_CANDIDATES", "80")))
    candidates: list[tuple[float, np.ndarray, tuple[float, float, float]]] = []
    try:
        if xtc.is_file():
            import mdtraj as md

            frame_index = 0
            last: tuple[float, np.ndarray, tuple[float, float, float]] | None = None
            for chunk in md.iterload(str(xtc), top=str(gro), chunk=25):
                for local in range(chunk.n_frames):
                    coords = np.asarray(chunk.xyz[local], dtype=float)
                    box = tuple(float(x) for x in chunk.unitcell_lengths[local][:3])
                    time_ps = float(chunk.time[local]) if chunk.time is not None and len(chunk.time) > local else float("nan")
                    last = (time_ps, coords.copy(), box)
                    if frame_index % candidate_stride == 0:
                        candidates.append(last)
                        if len(candidates) > max_candidates:
                            candidates = candidates[::2]
                            candidate_stride *= 2
                    frame_index += 1
            if last is not None and (not candidates or candidates[-1][0] != last[0]):
                candidates.append(last)
        if not candidates:
            candidates = [(float("nan"), gro_coords, gro_box)]
    except Exception:
        candidates = [(float("nan"), gro_coords, gro_box)]
    if len(candidates) > max_frames:
        keep = np.unique(np.linspace(0, len(candidates) - 1, max_frames, dtype=int))
        candidates = [candidates[int(i)] for i in keep]

    frame_paths: list[Path] = []
    for idx, (time_ps, coords, box) in enumerate(candidates):
        frame = case_dir / f"frame_{idx:03d}.png"
        out = _plot_structure_frame(
            species=species,
            coords=coords,
            box=box,
            item=item,
            out_png=frame,
            title=f"{item.label} trajectory overview",
            time_ps=time_ps,
        )
        if out is not None:
            frame_paths.append(out)
    if not frame_paths:
        return None, None
    poster = frame_paths[0]
    try:
        import imageio.v2 as imageio

        with imageio.get_writer(mp4, fps=2, codec="libx264", quality=6, macro_block_size=16) as writer:
            for frame in frame_paths:
                writer.append_data(imageio.imread(frame))
    except Exception:
        ffmpeg = shutil.which("ffmpeg")
        if ffmpeg:
            try:
                # Keep the video small and broadly compatible.  The frame list
                # is intentionally sparse; this movie is for visual t=0 ->
                # final sanity, not quantitative analysis.
                cmd = [
                    ffmpeg,
                    "-y",
                    "-hide_banner",
                    "-loglevel",
                    "error",
                    "-framerate",
                    "2",
                    "-i",
                    str(case_dir / "frame_%03d.png"),
                    "-vf",
                    "scale=trunc(iw/2)*2:trunc(ih/2)*2",
                    "-c:v",
                    "libx264",
                    "-pix_fmt",
                    "yuv420p",
                    "-crf",
                    "28",
                    str(mp4),
                ]
                subprocess.run(cmd, check=True)
            except Exception:
                return None, poster
        else:
            return None, poster
    return (mp4 if mp4.is_file() else None), poster


def _data_availability_rows(payloads: list[CasePayload], fig_dir: Path) -> list[dict[str, Any]]:
    rows: list[dict[str, Any]] = []
    for item in payloads:
        orientation_rows = _read_rows(fig_dir / "carbonyl_orientation_summary.csv")
        orientation_n = sum(int(_safe_float(row.get("sample_count"), 0.0)) for row in orientation_rows if row.get("case") == item.label)
        rdf_rows = _read_rows(item.analysis_dir / "time_series" / "rdf_cn_curves_timeseries.csv")
        ts_dir = item.analysis_dir / "time_series"
        rows.append(
            {
                "case": item.label,
                "charge_uC_cm2": item.charge,
                "z_density_rows": len(_read_rows(fig_dir / "zrel_csv" / f"{item.dirname}_z_density_profiles_zrel.csv")),
                "edl_rows": len(_read_rows(fig_dir / "zrel_csv" / f"{item.dirname}_electrostatic_potential_zrel.csv")),
                "membrane_timeseries_rows": len([row for row in _read_rows(fig_dir / "membrane_fraction_1ns_3point_smooth.csv") if row.get("case") == item.label]),
                "carbonyl_orientation_samples": orientation_n,
                "rdf_cn_rows": len(rdf_rows),
                "interface_rdf_cn_mp4": (ts_dir / "rdf_cn_timeseries.mp4").is_file(),
                "trajectory_overview_mp4": (fig_dir / "trajectory_overview" / f"{item.dirname}_trajectory_overview.mp4").is_file(),
                "final_xtc": (_final_nvt_dir(item) / "md.xtc").is_file(),
            }
        )
    _write_rows(fig_dir / "report_data_availability.csv", rows)
    return rows


def main() -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    _configure_matplotlib()
    root = Path(os.environ.get("EG08_SWEEP_ROOT", ".")).resolve()
    out_dir = root / "99_report"
    fig_dir = out_dir / "ppt_figures"
    out_dir.mkdir(parents=True, exist_ok=True)
    fig_dir.mkdir(parents=True, exist_ok=True)
    payloads = _load_payloads(root, _case_specs())
    _write_z_axis_reference(payloads, out_dir)

    t0_validation = _validate_t0(payloads, out_dir)
    z_arrays = _aggregate_z_profiles(payloads, fig_dir)
    z_by_charge = _plot_z_facets_by_charge(z_arrays, payloads, fig_dir)
    z_by_charge_zoom = _plot_z_facets_by_charge_zoom(z_arrays, payloads, fig_dir)
    z_by_species = _plot_z_facets_by_species(z_arrays, payloads, fig_dir)
    edl = _plot_edl_overlays(payloads, fig_dir)
    mem_data = _membrane_timeseries(payloads, fig_dir)
    mem_by_species, mem_by_charge = _plot_membrane_fraction(mem_data, payloads, fig_dir)
    penetration_metrics = _penetration_metrics(payloads, fig_dir)
    penetration_plots = _plot_penetration_metrics(penetration_metrics, payloads, fig_dir)
    penetration_schematic = _draw_penetration_schematic(fig_dir)
    orientation_summary, orientation_samples, orientation_plots = _orientation_rows(payloads, fig_dir)
    structure_plots = _structure_assets(payloads, fig_dir)
    trajectory_overviews = {item.label: _trajectory_assets(item, fig_dir) for item in payloads}
    availability_rows = _data_availability_rows(payloads, fig_dir)
    li_solvation_summary, li_solvation_plots = _li_solvation_depth_profiles(payloads, fig_dir)
    rdf_plots = _plot_rdf_cn_last_window(payloads, fig_dir)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Eg08.07 CMC-facing charge sweep"
    slide.placeholders[1].text = f"Strict post-processing report\nroot: {root}\nvalid_shared_t0={t0_validation.get('ok')}"

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Shared t=0 Validation"
    rows = [["case", "02_system/system.gro md5", "shared t0 OK"]]
    ref_md5 = None
    for row in t0_validation.get("system_gro_md5") or []:
        ref_md5 = ref_md5 or row.get("md5")
        rows.append([row.get("case"), row.get("md5"), bool(row.get("md5") == ref_md5)])
    if len(rows) == 1:
        rows.append(["not available", t0_validation.get("reason"), False])
    _add_table(slide, rows, Inches(0.45), Inches(0.85), Inches(12.4), Inches(1.7), font_size=8)
    traj = t0_validation.get("trajectory_first_frame") or {}
    traj_rows = [["trajectory first-frame diagnostic", "value"]]
    traj_rows.append(["available", traj.get("available")])
    traj_rows.append(["identical first saved xtc frame", traj.get("ok")])
    if traj.get("comparisons"):
        worst = max(float(row.get("first_frame_rmsd_nm") or 0.0) for row in traj.get("comparisons") or [])
        traj_rows.append(["worst first-frame RMSD / nm", f"{worst:.3e}"])
    else:
        traj_rows.append(["reason/note", traj.get("reason") or traj.get("note") or "not checked"])
    _add_table(slide, traj_rows, Inches(0.45), Inches(2.75), Inches(5.6), Inches(1.1), font_size=8)
    _add_picture(slide, structure_plots.get("shared_t0"), Inches(6.25), Inches(2.72), width=Inches(6.6))

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Case Data Availability"
    rows = [["case", "z rows", "EDL rows", "mem frac rows", "C=O samples", "RDF/CN rows", "traj MP4"]]
    for row in availability_rows:
        rows.append(
            [
                row.get("case"),
                row.get("z_density_rows"),
                row.get("edl_rows"),
                row.get("membrane_timeseries_rows"),
                row.get("carbonyl_orientation_samples"),
                row.get("rdf_cn_rows"),
                row.get("trajectory_overview_mp4"),
            ]
        )
    _add_table(slide, rows, Inches(0.45), Inches(0.85), Inches(12.4), Inches(1.7), font_size=8)
    _add_textbox(
        slide,
        Inches(0.65),
        Inches(2.95),
        Inches(12.0),
        Inches(3.25),
        "Purpose: this table separates missing data from real zero-valued trends. "
        "If a case has nonzero source rows, the PPT is allowed to plot that case even if an earlier run/root did not have it.\n"
        "Observation: the active report root and case list are shown on the title slide and in charge_sweep_full_ppt_paths.json.\n"
        "Analysis: blank facets should only occur with an explicit low-sample or missing-data note. A silent blank panel means the report script should be fixed, not interpreted physically.\n"
        "Conclusion: use this slide first when checking whether 0 or -9 uC/cm2 are absent, truly zero, or simply unavailable for one statistic.",
        font_size=13,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Unified z-axis Reference"
    z_rows = [["case", "z0/nm", "Lz/nm", "direction", "surface charge", "source"]]
    for item in payloads:
        z_rows.append(
            [
                item.label,
                f"{item.z0_nm:.4f}",
                f"{item.box_z_nm:.4f}",
                item.z_axis_direction,
                "" if item.z_axis_surface_charge_uC_cm2 is None else f"{item.z_axis_surface_charge_uC_cm2:.2f}",
                item.z_axis_source,
            ]
        )
    _add_table(slide, z_rows, Inches(0.45), Inches(0.85), Inches(12.4), Inches(1.85), font_size=7)
    _add_textbox(
        slide,
        Inches(0.65),
        Inches(3.05),
        Inches(12.0),
        Inches(2.8),
        "Definition: z_plot = 0 is the CMC-facing graphite inner surface, i.e. the surface that carries the requested negative charge in the sweep.\n"
        "Direction: +z_plot points from that graphite surface into CMC, then electrolyte, then the opposite graphite, and is wrapped into one complete periodic box [0, Lz).\n"
        "Acceptance rule: all z distribution, EDL charge/potential, penetration, Li solvation-depth, and z time-series plots must use this z_plot coordinate and must not show negative z axes.",
        font_size=14,
    )
    status = "PASS" if t0_validation.get("ok") else "FAIL"
    _add_textbox(
        slide,
        Inches(6.35),
        Inches(2.75),
        Inches(6.35),
        Inches(2.7),
        f"Observation: shared t=0 validation status is {status}.\n"
        "Analysis: the strict criterion is identical 02_system/system.gro files, because this is the shared assembled structure before charge-specific production. The first saved xtc frame is only a diagnostic; it may be after several integration steps.\n"
        "Conclusion: if the system.gro md5 check fails, all following physical trends are diagnostic-only and the four cases must be rerun from one shared_t0 structure.",
        font_size=14,
    )

    _result_slide(
        prs,
        "EDL Charge Density Overlay",
        edl.get("edl_charge_density_zrel"),
        "Observation: all charge states are plotted on the same z_plot axis.\nAnalysis: charge density localizes electrode and ionic charge relative to the negatively charged CMC-facing graphite surface.\nConclusion: use sign and peak position to identify screening/overscreening; negative x coordinates indicate a failed report.",
        source_svg=(fig_dir / "edl_charge_density_zrel.svg"),
    )
    _result_slide(
        prs,
        "EDL Integrated Charge Overlay",
        edl.get("edl_integrated_charge_zrel"),
        "Observation: integrated charge is recomputed cumulatively from z_plot=0.\nAnalysis: plateaus and sign reversals show how the CMC/electrolyte compensates the negative graphite surface charge.\nConclusion: compare these curves only after t=0 validation passes.",
        source_svg=(fig_dir / "edl_integrated_charge_zrel.svg"),
    )
    _result_slide(
        prs,
        "EDL Electrostatic Potential Overlay",
        edl.get("edl_potential_zrel"),
        "Observation: potential is obtained by integrating the reordered charge-derived electric field and then referencing to a local interior window.\nAnalysis: the absolute zero is arbitrary, but the negative surface should be read relative to the adjacent interior, not from a raw centered-z curve.\nConclusion: stronger CMC-facing negative charge should reshape the potential near the CMC-side EDL.",
        source_svg=(fig_dir / "edl_potential_zrel.svg"),
    )

    _result_slide(
        prs,
        "z Distribution: Charge Facets",
        z_by_charge["png"],
        "Observation: each subplot is one surface charge and compares all species.\nAnalysis: step-filled profiles are z-bin histograms, avoiding misleading overlaid line spaghetti.\nConclusion: this view answers which species occupy each region under one charge state.",
        source_svg=z_by_charge["svg"],
    )
    _result_slide(
        prs,
        "z Distribution: Charge Facets, CMC-interface Zoom",
        z_by_charge_zoom["png"],
        "Observation: this view magnifies the first few nanometers from the CMC-facing graphite surface with y fixed to 0-0.5 g cm-3.\n"
        "Analysis: it is designed to reveal dilute penetration tails and near-interface depletion that are invisible on the full-density scale.\n"
        "Conclusion: use this panel to judge early CMC-side infiltration and sparse species accumulation near the charged surface.",
        source_svg=z_by_charge_zoom["svg"],
    )
    _result_slide(
        prs,
        "z Distribution: Species Facets",
        z_by_species["png"],
        "Observation: each subplot is one species and overlays all charge states.\nAnalysis: shared z_plot axis makes charge-induced redistribution easier to read.\nConclusion: this view answers whether a given species responds monotonically to surface charge.",
        source_svg=z_by_species["svg"],
    )

    _result_slide(
        prs,
        "Membrane Permeation Metric Definitions",
        penetration_schematic["png"],
        "Definitions used throughout the report:\n"
        "f_mem(species,t)=N_membrane/(N_feed+N_membrane+N_permeate), a molecule-count fraction.\n"
        "P_entry=entry_event_count/initial_feed_count, a normalized entry-event index.\n"
        "D95 is the 95th percentile of penetration depth for molecule-frames that have entered CMC.\n"
        "AUC_depth integrates the normalized depth distribution and summarizes overall depth capability.\n"
        "entry flux=Delta entry_events/(Delta t x interface area), a rate-like comparison.",
        source_svg=penetration_schematic["svg"],
    )

    for sp, png in mem_by_species.items():
        _result_slide(
            prs,
            f"{sp} Membrane Fraction Across Charges",
            png,
            f"Definition: f_mem({sp},t)=N_membrane({sp},t)/(N_feed+N_membrane+N_permeate).\n"
            "Observation: curves use 1 ns bins with a 3-point rolling mean.\n"
            "Analysis: this is a molecule-count fraction, not mass fraction or permeability.\n"
            "Conclusion: sustained increases indicate membrane uptake; early differences require shared t=0 validation.",
            source_svg=png.with_suffix(".svg"),
        )
    for label, png in mem_by_charge.items():
        _result_slide(
            prs,
            f"{label} Membrane Fraction By Species",
            png,
            "Observation: all permeant species are compared within one charge state.\nAnalysis: solvent/cation/anion asymmetry reveals selective uptake.\nConclusion: this complements the per-species charge-sweep plots.",
            source_svg=png.with_suffix(".svg"),
        )

    for metric, png in penetration_plots.items():
        _result_slide(
            prs,
            f"Penetration Capability: {metric}",
            png,
            "Method: molecule COM is classified as feed, membrane, or permeate relative to the CMCNA membrane interval.\n"
            "Observation: bars compare EC/EMC/DEC/Li/PF6 within each surface charge.\n"
            "Analysis: P_entry is event-normalized, D95 describes the deep tail, AUC_depth summarizes the full normalized depth distribution, and loading reports mean membrane occupancy per volume.\n"
            "Conclusion: a species has stronger membrane-infiltration ability only when entry frequency and depth metrics increase together.",
            source_svg=png.with_suffix(".svg"),
        )

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Graphite EDL Carbonyl Orientation: Method And Samples"
    rows = [["case", "species", "n", "mean/deg", "median/deg", "IQR/deg", "status"]]
    for row in orientation_summary:
        n = int(_safe_float(row.get("sample_count"), 0.0))
        rows.append(
            [
                row.get("case"),
                row.get("species"),
                n,
                "" if n == 0 else f"{_safe_float(row.get('mean_angle_deg')):.1f}",
                "" if n == 0 else f"{_safe_float(row.get('median_angle_deg')):.1f}",
                "" if n == 0 else f"{_safe_float(row.get('q25_angle_deg')):.1f}-{_safe_float(row.get('q75_angle_deg')):.1f}",
                "low sample" if bool(row.get("low_sample")) else "ok",
            ]
        )
    _add_table(slide, rows[:15], Inches(0.45), Inches(0.85), Inches(12.4), Inches(3.4), font_size=7)
    _add_textbox(
        slide,
        Inches(0.6),
        Inches(4.65),
        Inches(12.0),
        Inches(1.7),
        "Only adsorbed EC/EMC/DEC molecules with orientation_available=True are included. "
        "The angle is carbonyl C=O relative to the local graphite surface normal: 0/180 deg is normal-like, 90 deg is parallel to graphite. "
        "Low-sample cells are diagnostic rather than physical trends.",
        font_size=12,
    )

    if orientation_plots.get("trend"):
        _result_slide(
            prs,
            "Graphite EDL Carbonyl Orientation vs Charge",
            orientation_plots["trend"],
            "Observation: EC/EMC/DEC mean carbonyl angle is plotted with IQR bars against CMC-facing surface charge.\n"
            "Analysis: systematic shifts toward 90 deg indicate carbonyls becoming more parallel to graphite; shifts toward 0/180 deg indicate stronger normal alignment.\n"
            "Conclusion: only interpret species/charges marked with sufficient oriented samples in the diagnostic table.",
            source_svg=orientation_plots["trend"].with_suffix(".svg"),
        )
    for sp in ("EC", "EMC", "DEC"):
        key = f"distribution_{sp}"
        if orientation_plots.get(key):
            _result_slide(
                prs,
                f"{sp} Carbonyl Orientation Distribution",
                orientation_plots[key],
                "Observation: each facet is one surface charge and the y-axis is the percent of adsorbed oriented frames.\n"
                "Analysis: broad or bimodal distributions mean multiple adsorption geometries coexist in the EDL.\n"
                "Conclusion: compare this distribution with the mean/IQR trend before making a charge-dependent orientation claim.",
                source_svg=orientation_plots[key].with_suffix(".svg"),
            )

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Li Solvation vs CMC Penetration Depth: Method"
    rows = [["case", "available", "Li-in-CMC samples", "stride", "cutoff / nm", "feed side"]]
    for row in li_solvation_summary:
        rows.append(
            [
                row.get("case"),
                row.get("available"),
                row.get("li_inside_cmc_frame_samples", row.get("reason")),
                row.get("frame_stride", ""),
                row.get("cutoff_nm", ""),
                row.get("feed_side", ""),
            ]
        )
    _add_table(slide, rows[:8], Inches(0.45), Inches(0.82), Inches(12.4), Inches(2.0), font_size=8)
    _add_textbox(
        slide,
        Inches(0.6),
        Inches(3.15),
        Inches(12.0),
        Inches(3.2),
        "Definition: only Li atoms located inside the CMC membrane interval are counted. "
        "Penetration depth is measured from the electrolyte-side CMC boundary toward the CMC interior.\n"
        "Solvation structure: for each depth bin, average the number of solvent representative O sites, CMC oxygen-like sites, and PF6 fluorine sites within the cutoff around Li. "
        "The plotted composition fractions are each component divided by solvent-O + CMC-O + PF6-F coordination.\n"
        "Interpretation: decreasing solvent-O with increasing CMC-O means Li sheds carbonate solvation and coordinates to the CMC matrix; PF6-F growth indicates ion-pair participation inside the membrane.",
        font_size=13,
    )

    for key, title in (
        ("cn_depth", "Li Solvation CN vs CMC Penetration Depth"),
        ("cn_site_charge_facets", "Li Solvation CN: Site Facets Across Charges"),
        ("fraction_depth", "Li Solvation Composition vs Depth"),
    ):
        png = li_solvation_plots.get(key)
        _result_slide(
            prs,
            title,
            png,
            "Observation: curves are binned by Li penetration depth inside CMC, not by simulation time.\n"
            "Analysis: CN values are local coordination counts around Li within the stated cutoff; values above four can occur when broad site classes are counted.\n"
            "Conclusion: this plot directly shows how the average Li solvation shell changes after entering the CMC membrane.",
            source_svg=(png.with_suffix(".svg") if png else None),
        )

    for pair, png in rdf_plots.items():
        _result_slide(
            prs,
            f"Interface RDF/CN: {pair}",
            png,
            "Observation: RDF is solid, CN is dashed on the right axis fixed to 0-6.\n"
            "Analysis: interface RDF uses local slab-aware normalization rather than bulk gmx rdf assumptions.\n"
            "Conclusion: first peak labels report preferred coordination distances; CN above 6 is recorded in CSV rather than stretching the plot.",
            source_svg=png.with_suffix(".svg"),
        )

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Structure Evolution Posters"
    positions = [(Inches(0.45), Inches(0.85)), (Inches(6.85), Inches(0.85)), (Inches(0.45), Inches(4.05)), (Inches(6.85), Inches(4.05))]
    for item, (left, top) in zip(payloads, positions):
        png = structure_plots.get(f"final_{item.dirname}")
        _add_picture(slide, png, left, top, width=Inches(5.85), height=Inches(2.48))
        _add_textbox(slide, left, top + Inches(2.52), Inches(5.85), Inches(0.35), f"{item.label}: final projection, same z_plot convention", font_size=8)

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Trajectory Overview MP4: Same t=0, Same View"
    for item, (left, top) in zip(payloads, positions):
        mp4, poster = trajectory_overviews.get(item.label, (None, None))
        _add_movie_tile(slide, mp4 or Path("missing.mp4"), poster, left, top, Inches(5.85), Inches(2.5))
        _add_textbox(slide, left, top + Inches(2.55), Inches(5.85), Inches(0.32), f"{item.label}: z_plot-y projection, sampled trajectory frames", font_size=8)

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Time-Series MP4 Panels"
    positions = [(Inches(0.45), Inches(0.85)), (Inches(6.85), Inches(0.85)), (Inches(0.45), Inches(4.05)), (Inches(6.85), Inches(4.05))]
    for item, (left, top) in zip(payloads, positions):
        mp4 = item.analysis_dir / "time_series" / "rdf_cn_timeseries.mp4"
        poster = item.analysis_dir / "time_series" / "frames" / "rdf_cn" / "frame_000.png"
        _add_movie_tile(slide, mp4, poster, left, top, Inches(5.85), Inches(2.5))
        _add_textbox(slide, left, top + Inches(2.55), Inches(5.85), Inches(0.30), f"{item.label}: RDF/CN time series", font_size=9)

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Source Data And Audit"
    _add_table(
        slide,
        [
            ["artifact", "path"],
            ["z_plot CSV", fig_dir / "zrel_csv"],
            ["z-axis reference", out_dir / "z_axis_reference.json"],
            ["membrane fraction CSV", fig_dir / "membrane_fraction_1ns_3point_smooth.csv"],
            ["penetration metrics CSV", fig_dir / "penetration_capability_metrics.csv"],
            ["Li solvation-depth CSV", fig_dir / "li_solvation_by_cmc_depth.csv"],
            ["z zoom CSV", fig_dir / "z_distribution_by_charge_facets_cmc_interface_zoom.csv"],
            ["carbonyl orientation summary", fig_dir / "carbonyl_orientation_summary.csv"],
            ["data availability CSV", fig_dir / "report_data_availability.csv"],
            ["trajectory overview MP4 directory", fig_dir / "trajectory_overview"],
            ["t0 validation JSON", out_dir / "shared_t0_validation.json"],
            ["z-axis audit JSON", out_dir / "z_axis_visual_audit.json"],
            ["visual audit JSON", out_dir / "ppt_visual_audit.json"],
        ],
        Inches(0.45),
        Inches(0.9),
        Inches(12.4),
        Inches(2.2),
        font_size=8,
    )
    _add_textbox(
        slide,
        Inches(0.6),
        Inches(3.55),
        Inches(12.0),
        Inches(2.4),
        "Observation: all generated figures have SVG sources and PNG previews for PowerPoint compatibility.\n"
        "Analysis: PPT itself is a reading layer; final plots can be redrawn from the listed CSV/SVG files.\n"
        "Conclusion: do not accept the PPT unless visual audit passes and file size remains below 50 MB.",
        font_size=14,
    )

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = "Arial"
                        if run.font.size is None:
                            run.font.size = Pt(12)

    out_ppt = out_dir / "eg08_07_charge_sweep_full_report.pptx"
    prs.save(out_ppt)
    audit = _audit_ppt(out_ppt, out_dir / "ppt_visual_audit.json")
    z_audit = _audit_z_axis(payloads, fig_dir, out_dir / "z_axis_visual_audit.json")
    artifact_checks = {
        "z_distribution_zoom_exists": (fig_dir / "z_distribution_by_charge_facets_cmc_interface_zoom.png").is_file(),
        "membrane_fraction_adaptive_exists": any((fig_dir / f"membrane_fraction_species_{sp}.png").is_file() for sp in PERMEANT_SPECIES),
        "penetration_schematic_exists": (fig_dir / "penetration_metric_schematic.png").is_file(),
        "carbonyl_orientation_or_diagnostic_exists": (fig_dir / "carbonyl_orientation_summary.csv").is_file(),
        "shared_t0_structure_exists": (fig_dir / "structure_overview" / "shared_t0_structure.png").is_file(),
    }
    audit["z_axis_ok"] = bool(z_audit["ok"])
    audit["artifact_checks"] = artifact_checks
    audit["ok"] = bool(audit["ok"] and z_audit["ok"] and all(artifact_checks.values()))
    _write_json(out_dir / "ppt_visual_audit.json", audit)
    final = {
        "ppt": str(out_ppt),
        "figure_dir": str(fig_dir),
        "root": str(root),
        "case_count": len(payloads),
        "valid_shared_t0": bool(t0_validation.get("ok")),
        "ppt_size_mb": audit["size_mb"],
        "visual_audit_ok": bool(audit["ok"]),
        "z_axis_audit_ok": bool(z_audit["ok"]),
        "li_solvation_depth_csv": str(fig_dir / "li_solvation_by_cmc_depth.csv"),
        "z_distribution_zoom_csv": str(fig_dir / "z_distribution_by_charge_facets_cmc_interface_zoom.csv"),
        "carbonyl_orientation_summary_csv": str(fig_dir / "carbonyl_orientation_summary.csv"),
        "penetration_schematic_svg": str(fig_dir / "penetration_metric_schematic.svg"),
    }
    _write_json(out_dir / "charge_sweep_full_ppt_paths.json", final)
    print(json.dumps(final, indent=2, ensure_ascii=False))


if __name__ == "__main__":
    main()
